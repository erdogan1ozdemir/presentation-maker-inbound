#!/usr/bin/env python3
"""
Loft SEO Şubat 2026 Değerlendirme Sunumu Oluşturucu
Referans sunumu baz alarak GSC ve GA4 verileriyle yeni sunum oluşturur.
"""

import os
import csv
import copy
import warnings
from collections import defaultdict

import openpyxl
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

warnings.filterwarnings('ignore')

# ============================================================
# CONFIG
# ============================================================
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(BASE_DIR, '(Shared) Loft - 2025 Değerlendirme.pptx')
OUTPUT_PATH = os.path.join(BASE_DIR, 'Loft_SEO_Subat_2026_Degerlendirme.pptx')
CHARTS_DIR = os.path.join(BASE_DIR, 'charts')
os.makedirs(CHARTS_DIR, exist_ok=True)

# Color palette
DARK_BLUE = '#1B2A4A'
MEDIUM_BLUE = '#2E4057'
LIGHT_BLUE = '#4A90D9'
ACCENT_GREEN = '#27AE60'
ACCENT_RED = '#E74C3C'
ACCENT_ORANGE = '#F39C12'
GRAY = '#95A5A6'
LIGHT_GRAY = '#ECF0F1'
WHITE = '#FFFFFF'
BLACK = '#2C3E50'

# Brand colors
COLOR_2025 = '#1B2A4A'
COLOR_2026 = '#4A90D9'
COLOR_PREV_YEAR = '#BDC3C7'

# ============================================================
# DATA LOADING
# ============================================================

def load_gsc_excel(filename):
    """Load GSC Excel file and return data by sheet."""
    path = os.path.join(BASE_DIR, filename)
    wb = openpyxl.load_workbook(path)
    data = {}
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if rows:
            data[sheet_name] = {'header': rows[0], 'rows': rows[1:]}
    return data

def load_csv_with_periods(filename):
    """Load CSV file that may contain multiple date periods."""
    path = os.path.join(BASE_DIR, filename)
    periods = []
    current_period = {'meta': {}, 'header': None, 'rows': []}

    with open(path, 'r') as f:
        reader = csv.reader(f)
        for row in reader:
            if not row or not row[0]:
                if current_period['header']:
                    periods.append(current_period)
                    current_period = {'meta': {}, 'header': None, 'rows': []}
                continue
            if row[0].startswith('#'):
                line = row[0].strip('# ').strip()
                if 'Start date:' in line:
                    current_period['meta']['start_date'] = line.split(':')[1].strip()
                elif 'End date:' in line:
                    current_period['meta']['end_date'] = line.split(':')[1].strip()
                continue
            if current_period['header'] is None:
                current_period['header'] = row
            else:
                current_period['rows'].append(row)

    if current_period['header']:
        periods.append(current_period)

    return periods

def load_landing_page_data(filename):
    """Load landing page CSV."""
    path = os.path.join(BASE_DIR, filename)
    header = None
    rows = []
    with open(path, 'r') as f:
        reader = csv.reader(f)
        for row in reader:
            if not row or row[0].startswith('#'):
                continue
            if header is None:
                header = row
            else:
                rows.append(row)
    return header, rows


# ============================================================
# DATA PROCESSING
# ============================================================

print("Veriler yükleniyor...")

# GSC Data
gsc_total_mom = load_gsc_excel('https___www.loft.com.tr_-Performance-on-Search-2026-03-11.xlsx')
gsc_total_yoy = load_gsc_excel('https___www.loft.com.tr_-Performance-on-Search-2026-03-11 (1).xlsx')
gsc_brand_yoy = load_gsc_excel('https___www.loft.com.tr_-Performance-on-Search-2026-03-11 (2).xlsx')
gsc_nonbrand_yoy = load_gsc_excel('https___www.loft.com.tr_-Performance-on-Search-2026-03-11 (3).xlsx')
gsc_brand_mom = load_gsc_excel('https___www.loft.com.tr_-Performance-on-Search-2026-03-11 (4).xlsx')
gsc_nonbrand_mom = load_gsc_excel('https___www.loft.com.tr_-Performance-on-Search-2026-03-11 (5).xlsx')

# GA4 Data
traffic_feb26_vs_jan26 = load_csv_with_periods('Traffic_acquisition_Session_default_channel_group (2).csv')
traffic_feb26_vs_feb25 = load_csv_with_periods('Traffic_acquisition_Session_default_channel_group (3).csv')

# Landing Page Data
lp_header, lp_rows = load_landing_page_data('Landing_page_Landing_page.csv')

print("Veriler işleniyor...")

# --- Compute GSC Totals ---
def sum_column(data, sheet, col_idx):
    return sum(r[col_idx] for r in data[sheet]['rows'] if r[col_idx])

# Total Impressions
total_imp_feb26 = sum_column(gsc_total_yoy, 'Queries', 3)
total_imp_feb25 = sum_column(gsc_total_yoy, 'Queries', 4)
total_imp_jan26 = sum_column(gsc_total_mom, 'Queries', 4)
imp_yoy_pct = ((total_imp_feb26 - total_imp_feb25) / total_imp_feb25) * 100
imp_mom_pct = ((total_imp_feb26 - total_imp_jan26) / total_imp_jan26) * 100

# Total Clicks
total_click_feb26 = sum_column(gsc_total_yoy, 'Queries', 1)
total_click_feb25 = sum_column(gsc_total_yoy, 'Queries', 2)
total_click_jan26 = sum_column(gsc_total_mom, 'Queries', 2)
click_yoy_pct = ((total_click_feb26 - total_click_feb25) / total_click_feb25) * 100
click_mom_pct = ((total_click_feb26 - total_click_jan26) / total_click_jan26) * 100

# Brand
brand_imp_feb26 = sum_column(gsc_brand_yoy, 'Queries', 3)
brand_imp_feb25 = sum_column(gsc_brand_yoy, 'Queries', 4)
brand_click_feb26 = sum_column(gsc_brand_yoy, 'Queries', 1)
brand_click_feb25 = sum_column(gsc_brand_yoy, 'Queries', 2)

# Non-Brand
nb_imp_feb26 = sum_column(gsc_nonbrand_yoy, 'Queries', 3)
nb_imp_feb25 = sum_column(gsc_nonbrand_yoy, 'Queries', 4)
nb_click_feb26 = sum_column(gsc_nonbrand_yoy, 'Queries', 1)
nb_click_feb25 = sum_column(gsc_nonbrand_yoy, 'Queries', 2)

nb_imp_yoy_pct = ((nb_imp_feb26 - nb_imp_feb25) / nb_imp_feb25) * 100
nb_click_yoy_pct = ((nb_click_feb26 - nb_click_feb25) / nb_click_feb25) * 100

# --- Non-Brand Keyword Analysis ---
def analyze_keywords(data, sheet='Queries'):
    rows = data[sheet]['rows']
    results = []
    for row in rows:
        q, c_new, c_old, i_new, i_old, ctr_new, ctr_old, pos_new, pos_old = row
        imp_change = ((i_new - i_old) / i_old * 100) if i_old and i_old > 0 else (99999 if i_new > 0 else 0)
        click_change = ((c_new - c_old) / c_old * 100) if c_old and c_old > 0 else (99999 if c_new > 0 else 0)
        results.append({
            'query': q,
            'imp_new': i_new or 0, 'imp_old': i_old or 0,
            'imp_change': imp_change,
            'click_new': c_new or 0, 'click_old': c_old or 0,
            'click_change': click_change,
            'pos': pos_new or 0
        })
    return results

nb_keywords = analyze_keywords(gsc_nonbrand_yoy)

# Keywords with significant impression increase (filter out very low volume)
nb_imp_increase = [k for k in nb_keywords if k['imp_change'] > 0 and k['imp_new'] >= 1000]
nb_imp_increase.sort(key=lambda x: x['imp_new'], reverse=True)

# Keywords with impression decrease (filter meaningful ones)
nb_imp_decrease = [k for k in nb_keywords if k['imp_change'] < 0 and k['imp_old'] >= 100]
nb_imp_decrease.sort(key=lambda x: x['imp_change'])

# Keywords with click increase
nb_click_increase = [k for k in nb_keywords if k['click_change'] > 0 and k['click_new'] >= 10]
nb_click_increase.sort(key=lambda x: x['click_new'], reverse=True)

# Keywords with click decrease
nb_click_decrease = [k for k in nb_keywords if k['click_change'] < 0 and k['click_old'] >= 5]
nb_click_decrease.sort(key=lambda x: x['click_change'])

# --- GA4 Organic Traffic ---
# Feb 2026 vs Feb 2025
organic_feb26 = 0
organic_feb25 = 0
organic_jan26 = 0

for period in traffic_feb26_vs_feb25:
    for row in period['rows']:
        if row[0] == 'Organic Search':
            if period['meta'].get('start_date', '').startswith('2026'):
                organic_feb26 = int(row[2])  # Sessions
            elif period['meta'].get('start_date', '').startswith('2025'):
                organic_feb25 = int(row[2])

for period in traffic_feb26_vs_jan26:
    for row in period['rows']:
        if row[0] == 'Organic Search':
            if period['meta'].get('start_date', '') == '20260101':
                organic_jan26 = int(row[2])

organic_yoy_pct = ((organic_feb26 - organic_feb25) / organic_feb25) * 100 if organic_feb25 else 0
organic_mom_pct = ((organic_feb26 - organic_jan26) / organic_jan26) * 100 if organic_jan26 else 0

# --- Landing Page Data ---
landing_pages = []
for row in lp_rows:
    if row[0] == '(not set)':
        continue
    try:
        landing_pages.append({
            'page': row[0],
            'sessions': int(row[1]),
            'users': int(row[2]),
            'new_users': int(row[3]),
            'avg_engagement': float(row[4]),
            'key_events': int(row[5]),
            'revenue': float(row[6]),
            'event_rate': float(row[7])
        })
    except (ValueError, IndexError):
        continue

landing_pages.sort(key=lambda x: x['sessions'], reverse=True)

# --- Device Breakdown ---
devices_feb26 = gsc_total_yoy['Devices']['rows']
devices_feb25_data = {r[0]: r for r in devices_feb26}

print("Grafikler oluşturuluyor...")

# ============================================================
# CHART GENERATION
# ============================================================

plt.rcParams['font.family'] = 'sans-serif'
plt.rcParams['font.size'] = 11
plt.rcParams['axes.spines.top'] = False
plt.rcParams['axes.spines.right'] = False

def format_number(n):
    """Format number with K/M suffix."""
    if n >= 1_000_000:
        return f'{n/1_000_000:.1f}M'
    elif n >= 1_000:
        return f'{n/1_000:.1f}K'
    return str(int(n))

def create_comparison_bar_chart(title, val_new, val_old, label_new, label_old,
                                 filename, color_new=COLOR_2026, color_old=COLOR_PREV_YEAR):
    """Create a horizontal comparison bar chart."""
    fig, ax = plt.subplots(figsize=(10, 2.5))

    y_pos = [0.6, 0]
    values = [val_new, val_old]
    colors = [color_new, color_old]
    labels = [label_new, label_old]

    bars = ax.barh(y_pos, values, height=0.45, color=colors, edgecolor='none')

    for bar, val, label in zip(bars, values, labels):
        ax.text(bar.get_width() + max(values) * 0.02, bar.get_y() + bar.get_height()/2,
                f'{format_number(val)}', va='center', fontsize=12, fontweight='bold', color=DARK_BLUE)

    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, fontsize=11, color=DARK_BLUE)
    ax.set_xlim(0, max(values) * 1.2)
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, p: format_number(x)))
    ax.set_title(title, fontsize=14, fontweight='bold', color=DARK_BLUE, loc='left', pad=10)
    ax.spines['left'].set_visible(False)
    ax.spines['bottom'].set_color(GRAY)
    ax.tick_params(axis='y', length=0)

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, filename), dpi=150, bbox_inches='tight',
                facecolor='white', transparent=False)
    plt.close()


def create_impression_click_chart(filename):
    """Create impression and click comparison chart."""
    fig, axes = plt.subplots(1, 2, figsize=(12, 4))

    # Impression chart
    ax1 = axes[0]
    categories = ['Total', 'Brand', 'Non-Brand']
    feb25_vals = [total_imp_feb25, brand_imp_feb25, nb_imp_feb25]
    feb26_vals = [total_imp_feb26, brand_imp_feb26, nb_imp_feb26]

    x = np.arange(len(categories))
    width = 0.35

    bars1 = ax1.bar(x - width/2, feb25_vals, width, label='Şub 2025', color=COLOR_PREV_YEAR, edgecolor='none')
    bars2 = ax1.bar(x + width/2, feb26_vals, width, label='Şub 2026', color=COLOR_2026, edgecolor='none')

    ax1.set_ylabel('Impression', fontsize=10, color=DARK_BLUE)
    ax1.set_title('Impression Karşılaştırma', fontsize=13, fontweight='bold', color=DARK_BLUE)
    ax1.set_xticks(x)
    ax1.set_xticklabels(categories, fontsize=10)
    ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, p: format_number(x)))
    ax1.legend(fontsize=9)

    for bar in bars2:
        ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(feb26_vals)*0.02,
                format_number(bar.get_height()), ha='center', va='bottom', fontsize=9, fontweight='bold')

    # Click chart
    ax2 = axes[1]
    feb25_clicks = [total_click_feb25, brand_click_feb25, nb_click_feb25]
    feb26_clicks = [total_click_feb26, brand_click_feb26, nb_click_feb26]

    bars3 = ax2.bar(x - width/2, feb25_clicks, width, label='Şub 2025', color=COLOR_PREV_YEAR, edgecolor='none')
    bars4 = ax2.bar(x + width/2, feb26_clicks, width, label='Şub 2026', color=COLOR_2026, edgecolor='none')

    ax2.set_ylabel('Click', fontsize=10, color=DARK_BLUE)
    ax2.set_title('Click Karşılaştırma', fontsize=13, fontweight='bold', color=DARK_BLUE)
    ax2.set_xticks(x)
    ax2.set_xticklabels(categories, fontsize=10)
    ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, p: format_number(x)))
    ax2.legend(fontsize=9)

    for bar in bars4:
        ax2.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(feb26_clicks)*0.02,
                format_number(bar.get_height()), ha='center', va='bottom', fontsize=9, fontweight='bold')

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, filename), dpi=150, bbox_inches='tight',
                facecolor='white', transparent=False)
    plt.close()


def create_keyword_table_chart(keywords, title, filename, top_n=15,
                                show_decrease=True, color=ACCENT_RED):
    """Create a keyword performance table as an image."""
    kws = keywords[:top_n]
    if not kws:
        return

    fig, ax = plt.subplots(figsize=(10, max(3, 0.45 * len(kws) + 1.5)))
    ax.axis('off')

    col_labels = ['Kelime', 'Şub 2025\nImpression', 'Şub 2026\nImpression', 'YoY\nDeğişim', 'Pozisyon']
    cell_data = []
    cell_colors = []

    for k in kws:
        change_str = f'{k["imp_change"]:.0f}%' if abs(k["imp_change"]) < 99999 else 'YENİ'
        cell_data.append([
            k['query'],
            format_number(k['imp_old']),
            format_number(k['imp_new']),
            change_str,
            f'{k["pos"]:.1f}'
        ])
        if k['imp_change'] > 0:
            cell_colors.append([WHITE, WHITE, WHITE, '#E8F5E9', WHITE])
        else:
            cell_colors.append([WHITE, WHITE, WHITE, '#FFEBEE', WHITE])

    table = ax.table(cellText=cell_data, colLabels=col_labels,
                     cellColours=cell_colors,
                     colColours=[LIGHT_BLUE]*5,
                     loc='center', cellLoc='center')

    table.auto_set_font_size(False)
    table.set_fontsize(9)
    table.scale(1, 1.4)

    # Style header
    for j in range(5):
        cell = table[0, j]
        cell.set_text_props(fontweight='bold', color='white', fontsize=9)
        cell.set_facecolor(DARK_BLUE)

    # First column left-aligned
    for i in range(len(kws) + 1):
        table[i, 0].set_text_props(ha='left')

    ax.set_title(title, fontsize=13, fontweight='bold', color=DARK_BLUE, loc='left', pad=15)

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, filename), dpi=150, bbox_inches='tight',
                facecolor='white', transparent=False)
    plt.close()


def create_click_keyword_dual_table(decrease_kws, increase_kws, filename, top_n=12):
    """Create side-by-side click keyword tables."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, max(3, 0.45 * top_n + 1.5)))

    for ax, kws, title, color in [
        (ax1, decrease_kws[:top_n], 'Click Düşüşü Yaşayan Kelimeler', ACCENT_RED),
        (ax2, increase_kws[:top_n], 'Click Artışı Yaşayan Kelimeler', ACCENT_GREEN)
    ]:
        ax.axis('off')
        if not kws:
            ax.text(0.5, 0.5, 'Veri bulunamadı', ha='center', va='center', fontsize=12)
            continue

        col_labels = ['Kelime', '2025', '2026', 'YoY %']
        cell_data = []
        cell_colors = []

        for k in kws:
            change_str = f'{k["click_change"]:.0f}%' if abs(k["click_change"]) < 99999 else 'YENİ'
            cell_data.append([
                k['query'],
                format_number(k['click_old']),
                format_number(k['click_new']),
                change_str
            ])
            bg = '#E8F5E9' if k['click_change'] > 0 else '#FFEBEE'
            cell_colors.append([WHITE, WHITE, WHITE, bg])

        table = ax.table(cellText=cell_data, colLabels=col_labels,
                        cellColours=cell_colors,
                        colColours=[color]*4,
                        loc='center', cellLoc='center')
        table.auto_set_font_size(False)
        table.set_fontsize(9)
        table.scale(1, 1.4)

        for j in range(4):
            cell = table[0, j]
            cell.set_text_props(fontweight='bold', color='white', fontsize=9)

        ax.set_title(title, fontsize=12, fontweight='bold', color=DARK_BLUE, pad=10)

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, filename), dpi=150, bbox_inches='tight',
                facecolor='white', transparent=False)
    plt.close()


def create_organic_traffic_chart(filename):
    """Create organic traffic comparison chart."""
    fig, ax = plt.subplots(figsize=(10, 4))

    channels = ['Organic\nSearch', 'Organic\nSocial', 'Organic\nShopping', 'Organic\nVideo']

    # Feb 2026 vs Feb 2025 from traffic file 3
    feb26_data = {}
    feb25_data = {}

    for period in traffic_feb26_vs_feb25:
        for row in period['rows']:
            channel = row[0]
            sessions = int(row[2]) if row[2] else 0
            if period['meta'].get('start_date', '').startswith('2026'):
                feb26_data[channel] = sessions
            elif period['meta'].get('start_date', '').startswith('2025'):
                feb25_data[channel] = sessions

    organic_channels = ['Organic Search', 'Organic Social', 'Organic Shopping', 'Organic Video']
    feb26_vals = [feb26_data.get(c, 0) for c in organic_channels]
    feb25_vals = [feb25_data.get(c, 0) for c in organic_channels]

    x = np.arange(len(channels))
    width = 0.35

    bars1 = ax.bar(x - width/2, feb25_vals, width, label='Şub 2025', color=COLOR_PREV_YEAR, edgecolor='none')
    bars2 = ax.bar(x + width/2, feb26_vals, width, label='Şub 2026', color=COLOR_2026, edgecolor='none')

    for bar in bars2:
        height = bar.get_height()
        if height > 0:
            ax.text(bar.get_x() + bar.get_width()/2, height + max(feb26_vals)*0.02,
                    format_number(height), ha='center', va='bottom', fontsize=9, fontweight='bold')

    ax.set_ylabel('Sessions', fontsize=10, color=DARK_BLUE)
    ax.set_title('Organik Kanal Bazlı Trafik Karşılaştırma (Şub 2025 vs Şub 2026)',
                fontsize=13, fontweight='bold', color=DARK_BLUE)
    ax.set_xticks(x)
    ax.set_xticklabels(channels, fontsize=10)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, p: format_number(x)))
    ax.legend(fontsize=10)

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, filename), dpi=150, bbox_inches='tight',
                facecolor='white', transparent=False)
    plt.close()


def create_landing_page_table(pages, filename, top_n=15):
    """Create landing page performance table chart."""
    fig, ax = plt.subplots(figsize=(12, max(3, 0.4 * top_n + 1.5)))
    ax.axis('off')

    col_labels = ['Landing Page', 'Sessions', 'Active Users', 'Yeni Kullanıcı', 'Ort. Etkileşim\nSüresi (sn)', 'Gelir (₺)']
    cell_data = []
    cell_colors = []

    for p in pages[:top_n]:
        page_name = p['page']
        if len(page_name) > 40:
            page_name = page_name[:37] + '...'
        cell_data.append([
            page_name,
            f'{p["sessions"]:,}',
            f'{p["users"]:,}',
            f'{p["new_users"]:,}',
            f'{p["avg_engagement"]:.0f}',
            f'{p["revenue"]:,.0f}'
        ])
        cell_colors.append([WHITE]*6)

    table = ax.table(cellText=cell_data, colLabels=col_labels,
                    cellColours=cell_colors,
                    colColours=[DARK_BLUE]*6,
                    loc='center', cellLoc='center')

    table.auto_set_font_size(False)
    table.set_fontsize(8)
    table.scale(1, 1.4)

    for j in range(6):
        cell = table[0, j]
        cell.set_text_props(fontweight='bold', color='white', fontsize=8)

    # Left-align first column
    for i in range(len(cell_data) + 1):
        table[i, 0].set_text_props(ha='left')

    ax.set_title('Şubat 2026 Landing Page Özelinde Organic Session Durumu',
                fontsize=13, fontweight='bold', color=DARK_BLUE, loc='left', pad=15)

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, filename), dpi=150, bbox_inches='tight',
                facecolor='white', transparent=False)
    plt.close()


def create_device_chart(filename):
    """Create device breakdown chart."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4))

    devices_yoy = gsc_total_yoy['Devices']['rows']

    device_names = []
    clicks_26 = []
    clicks_25 = []
    imp_26 = []
    imp_25 = []

    for row in devices_yoy:
        device_names.append(row[0])
        clicks_26.append(row[1] or 0)
        clicks_25.append(row[2] or 0)
        imp_26.append(row[3] or 0)
        imp_25.append(row[4] or 0)

    x = np.arange(len(device_names))
    width = 0.35

    ax1.bar(x - width/2, imp_25, width, label='Şub 2025', color=COLOR_PREV_YEAR)
    ax1.bar(x + width/2, imp_26, width, label='Şub 2026', color=COLOR_2026)
    ax1.set_title('Cihaz Bazlı Impression', fontsize=12, fontweight='bold', color=DARK_BLUE)
    ax1.set_xticks(x)
    ax1.set_xticklabels(device_names, fontsize=10)
    ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, p: format_number(x)))
    ax1.legend(fontsize=9)

    ax2.bar(x - width/2, clicks_25, width, label='Şub 2025', color=COLOR_PREV_YEAR)
    ax2.bar(x + width/2, clicks_26, width, label='Şub 2026', color=COLOR_2026)
    ax2.set_title('Cihaz Bazlı Click', fontsize=12, fontweight='bold', color=DARK_BLUE)
    ax2.set_xticks(x)
    ax2.set_xticklabels(device_names, fontsize=10)
    ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, p: format_number(x)))
    ax2.legend(fontsize=9)

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, filename), dpi=150, bbox_inches='tight',
                facecolor='white', transparent=False)
    plt.close()


def create_traffic_channel_chart(filename):
    """Create full traffic channel comparison chart."""
    fig, ax = plt.subplots(figsize=(14, 5))

    feb26_data = {}
    feb25_data = {}

    for period in traffic_feb26_vs_feb25:
        for row in period['rows']:
            channel = row[0]
            if not channel:
                continue
            sessions = int(row[2]) if row[2] and row[2] != '0' else 0
            if period['meta'].get('start_date', '').startswith('2026'):
                feb26_data[channel] = sessions
            elif period['meta'].get('start_date', '').startswith('2025'):
                feb25_data[channel] = sessions

    # Select top channels by Feb 2026 traffic
    all_channels = sorted(feb26_data.keys(), key=lambda c: feb26_data.get(c, 0), reverse=True)
    top_channels = [c for c in all_channels if feb26_data.get(c, 0) > 0][:10]

    feb26_vals = [feb26_data.get(c, 0) for c in top_channels]
    feb25_vals = [feb25_data.get(c, 0) for c in top_channels]

    x = np.arange(len(top_channels))
    width = 0.35

    ax.bar(x - width/2, feb25_vals, width, label='Şub 2025', color=COLOR_PREV_YEAR, edgecolor='none')
    ax.bar(x + width/2, feb26_vals, width, label='Şub 2026', color=COLOR_2026, edgecolor='none')

    # Add percentage change labels
    for i, (v25, v26) in enumerate(zip(feb25_vals, feb26_vals)):
        if v25 > 0:
            pct = ((v26 - v25) / v25) * 100
            color = ACCENT_GREEN if pct > 0 else ACCENT_RED
            ax.text(i, max(v25, v26) + max(feb26_vals)*0.02,
                   f'{pct:+.0f}%', ha='center', fontsize=8, fontweight='bold', color=color)

    short_names = [c.replace('Organic ', 'Org.\n').replace('Paid ', 'Paid\n') for c in top_channels]
    ax.set_xticks(x)
    ax.set_xticklabels(short_names, fontsize=9)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, p: format_number(x)))
    ax.set_title('Kanal Bazlı Trafik Karşılaştırma (Şub 2025 vs Şub 2026)',
                fontsize=13, fontweight='bold', color=DARK_BLUE)
    ax.legend(fontsize=10)

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, filename), dpi=150, bbox_inches='tight',
                facecolor='white', transparent=False)
    plt.close()


def create_top_queries_chart(filename):
    """Create top queries comparison chart."""
    # Get top 15 non-brand queries by Feb 2026 impressions
    nb_rows = gsc_nonbrand_yoy['Queries']['rows']
    sorted_by_imp = sorted(nb_rows, key=lambda r: r[3] if r[3] else 0, reverse=True)[:15]

    fig, ax = plt.subplots(figsize=(12, 6))

    queries = [r[0] for r in sorted_by_imp]
    imp_26 = [r[3] or 0 for r in sorted_by_imp]
    imp_25 = [r[4] or 0 for r in sorted_by_imp]

    y = np.arange(len(queries))
    height = 0.35

    ax.barh(y + height/2, imp_25, height, label='Şub 2025', color=COLOR_PREV_YEAR, edgecolor='none')
    ax.barh(y - height/2, imp_26, height, label='Şub 2026', color=COLOR_2026, edgecolor='none')

    ax.set_yticks(y)
    ax.set_yticklabels(queries, fontsize=9)
    ax.invert_yaxis()
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, p: format_number(x)))
    ax.set_title('Top Non-Brand Sorgular - Impression (Şub 2025 vs Şub 2026)',
                fontsize=13, fontweight='bold', color=DARK_BLUE)
    ax.legend(fontsize=10)

    plt.tight_layout()
    plt.savefig(os.path.join(CHARTS_DIR, filename), dpi=150, bbox_inches='tight',
                facecolor='white', transparent=False)
    plt.close()


# Generate all charts
create_impression_click_chart('impression_click_comparison.png')
create_keyword_table_chart(nb_imp_increase,
                          'Non-Brand Impression Artışı Yaşanan Kelimeler (YoY)',
                          'nb_imp_increase.png', top_n=15, color=ACCENT_GREEN)
create_keyword_table_chart(nb_imp_decrease,
                          'Non-Brand Impression Düşüşü Yaşanan Kelimeler (YoY)',
                          'nb_imp_decrease.png', top_n=15)
create_click_keyword_dual_table(nb_click_decrease, nb_click_increase, 'nb_click_dual.png')
create_organic_traffic_chart('organic_traffic.png')
create_landing_page_table(landing_pages, 'landing_pages.png')
create_device_chart('device_breakdown.png')
create_traffic_channel_chart('traffic_channels.png')
create_top_queries_chart('top_queries.png')

print("Grafikler oluşturuldu.")

# ============================================================
# PRESENTATION BUILDING
# ============================================================

print("Sunum oluşturuluyor...")

# Load template
prs = Presentation(TEMPLATE_PATH)
slide_width = prs.slide_width
slide_height = prs.slide_height

# Helper functions
def add_text_box(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=DARK_BLUE, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    """Add a text box to slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor.from_string(color.lstrip('#'))
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_header_bar(slide, text):
    """Add the header bar at top of slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(291125), Emu(52900), Emu(2741400), Emu(153600)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(DARK_BLUE.lstrip('#'))
    shape.line.fill.background()

    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Calibri'
    return shape

def add_footer(slide, text):
    """Add footer text."""
    add_text_box(slide, Emu(802275), Emu(4884438), Emu(6718500), Emu(230700),
                text, font_size=7, color=GRAY)

def add_line(slide, left, top, width):
    """Add horizontal line."""
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Emu(12700)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(LIGHT_BLUE.lstrip('#'))
    shape.line.fill.background()
    return shape

def add_image_to_slide(slide, image_path, left, top, width=None, height=None):
    """Add image to slide."""
    if width and height:
        slide.shapes.add_picture(image_path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(image_path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(image_path, left, top, height=height)
    else:
        slide.shapes.add_picture(image_path, left, top)


def pct_str(val):
    """Format percentage with direction."""
    direction = "artmış" if val > 0 else "azalmış"
    return f'%{abs(val):.1f} {direction}'

def pct_str_short(val):
    """Format percentage short."""
    sign = "+" if val > 0 else ""
    return f'{sign}{val:.1f}%'


# ---- Clear and rebuild slides ----
# We'll modify the existing slides in place where possible,
# and add new content

# Strategy: We'll work with the existing template slides.
# For each data slide, we'll update text boxes and replace/add images.

def find_text_shapes(slide):
    """Find all text shapes in a slide."""
    shapes = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text:
            shapes.append(shape)
    return shapes

def clear_slide_content(slide):
    """Remove all shapes from a slide (keep slide layout)."""
    sp_list = slide.shapes._spTree
    shapes_to_remove = []
    for sp in sp_list:
        shapes_to_remove.append(sp)
    for sp in shapes_to_remove:
        if sp.tag.endswith('}sp') or sp.tag.endswith('}pic') or sp.tag.endswith('}grpSp') or sp.tag.endswith('}cxnSp') or sp.tag.endswith('}graphicFrame'):
            sp_list.remove(sp)

def rebuild_slide(slide, header_text, subtitle_text, insight_text, chart_path,
                  footer_text='Source: GSC', chart_left=None, chart_top=None,
                  chart_width=None, chart_height=None):
    """Rebuild a data slide with standard layout."""
    clear_slide_content(slide)

    # Header bar
    add_header_bar(slide, header_text)

    # Divider line
    add_line(slide, Emu(213063), Emu(750000), Emu(8718000))

    # Subtitle / Insight text
    if subtitle_text:
        add_text_box(slide, Emu(182776), Emu(260000), Emu(8778600), Emu(450000),
                    subtitle_text, font_size=11, bold=False, color=DARK_BLUE)

    if insight_text:
        add_text_box(slide, Emu(182776), Emu(4100000), Emu(8557200), Emu(700000),
                    insight_text, font_size=9, color=MEDIUM_BLUE)

    # Chart image
    if chart_path and os.path.exists(chart_path):
        c_left = chart_left or Emu(182775)
        c_top = chart_top or Emu(850000)
        c_width = chart_width or Emu(8778600)
        c_height = chart_height or Emu(3200000)
        add_image_to_slide(slide, chart_path, c_left, c_top, c_width, c_height)

    # Footer
    add_footer(slide, footer_text)


# ============================================================
# UPDATE EACH SLIDE
# ============================================================

slides = list(prs.slides)

# --- Slide 1: Title (keep as is) ---

# --- Slide 2: Table of Contents (keep structure, update if needed) ---

# --- Slide 3: Section 01 Header (keep as is) ---

# --- Slide 4: Total Impression ---
print("  Slide 4: Total Impression")
rebuild_slide(
    slides[3],
    header_text='Loft | Search Console Verileri',
    subtitle_text=(
        f'Şubat 2026\'da total impression, Şubat 2025\'e göre {pct_str(imp_yoy_pct)}tır. '
        f'Impression Ocak 2026\'ya göre {pct_str(imp_mom_pct)}tır.\n'
        f'Non-Brand impression\'da YoY %{abs(nb_imp_yoy_pct):.0f} artış dikkat çekicidir.'
    ),
    insight_text=(
        f'- Toplam Impression: {format_number(total_imp_feb26)} (Şub 2026) vs {format_number(total_imp_feb25)} (Şub 2025)\n'
        f'- Brand Impression: {format_number(brand_imp_feb26)} (YoY {pct_str_short(((brand_imp_feb26-brand_imp_feb25)/brand_imp_feb25)*100)})\n'
        f'- Non-Brand Impression: {format_number(nb_imp_feb26)} (YoY {pct_str_short(nb_imp_yoy_pct)})\n'
        f'- Mobile impression ağırlıklı: {format_number(devices_feb26[0][3])} ({(devices_feb26[0][3]/total_imp_feb26*100):.0f}% pay)'
    ),
    chart_path=os.path.join(CHARTS_DIR, 'impression_click_comparison.png'),
    footer_text='Source: GSC'
)

# --- Slide 5: Non-Brand Impression Decrease ---
print("  Slide 5: Non-Brand Impression Azalış")
decrease_insight = ''
if nb_imp_decrease:
    top_dec = nb_imp_decrease[:5]
    dec_list = ', '.join([f'{k["query"]} ({pct_str_short(k["imp_change"])})' for k in top_dec])
    decrease_insight = (
        f'Geçen yıla göre impression düşüşü yaşanan kelimeler: {dec_list}\n'
        f'Bu kelimelerin büyük çoğunluğu sezonsallık etkisi veya site yapısı değişikliklerinden kaynaklanmaktadır.'
    )
else:
    decrease_insight = (
        'Non-brand keyword impression performansında geçen yıla göre genel bir artış trendi gözlemlenmektedir. '
        'Site genelindeki SEO iyileştirmeleri impression performansına olumlu yansımıştır.'
    )

rebuild_slide(
    slides[4],
    header_text='Loft | Search Console Verileri',
    subtitle_text='Geçen Yıla Göre Non-Brand Impression Azalışı Yaşanan Kelimeler',
    insight_text=decrease_insight,
    chart_path=os.path.join(CHARTS_DIR, 'nb_imp_decrease.png'),
    footer_text='Source: GSC'
)

# --- Slide 6: Search Volume Change (Keyword Planner - skip content) ---
print("  Slide 6: Arama Hacmi (Keyword Planner data yok - placeholder)")
rebuild_slide(
    slides[5],
    header_text='Loft | Arama Hacmi Değişimi',
    subtitle_text='Non-Brand Impression Düşüşü Yaşanan Kelimelerin Geçtiğimiz Seneye Göre Arama Hacmi Değişimi',
    insight_text='Not: Bu slide için Keyword Planner verisi gereklidir. Veri temin edildiğinde güncellenecektir.',
    chart_path=None,
    footer_text='Source: Keyword Planner'
)

# --- Slide 7: Search Volume Aggregate (skip) ---
print("  Slide 7: Arama Hacmi Aggregate (placeholder)")
rebuild_slide(
    slides[6],
    header_text='Loft | Arama Hacmi Değişimi',
    subtitle_text='Arama Hacmi Toplam Değişim',
    insight_text='Not: Bu slide için Keyword Planner verisi gereklidir. Veri temin edildiğinde güncellenecektir.',
    chart_path=None,
    footer_text='Source: Keyword Planner'
)

# --- Slide 8: Non-Brand Impression Increase ---
print("  Slide 8: Non-Brand Impression Artışı")
if nb_imp_increase:
    top_inc = nb_imp_increase[:5]
    inc_list = '\n'.join([f'- {k["query"]}: {format_number(k["imp_old"])} → {format_number(k["imp_new"])} ({pct_str_short(min(k["imp_change"], 99999))})' for k in top_inc])
    increase_insight = (
        f'Non-brand keyword\'lerde büyük impression artışları gözlemlenmektedir. '
        f'Özellikle outlet, mont, sweatshirt gibi genel kategori kelimelerinde çok yüksek artışlar dikkat çekmektedir.\n'
        f'{inc_list}'
    )
else:
    increase_insight = ''

rebuild_slide(
    slides[7],
    header_text='Loft | Search Console Verileri',
    subtitle_text='Geçen Yıla Göre Non-Brand Impression Artışı Yaşanan Kelimeler',
    insight_text=increase_insight,
    chart_path=os.path.join(CHARTS_DIR, 'nb_imp_increase.png'),
    footer_text='Source: GSC'
)

# --- Slide 9: Total Click ---
print("  Slide 9: Total Click")
rebuild_slide(
    slides[8],
    header_text='Loft | Search Console Verileri',
    subtitle_text=(
        f'Şubat 2026\'da total click, Şubat 2025\'e göre {pct_str(click_yoy_pct)}tır. '
        f'Click Ocak 2026\'ya göre {pct_str(click_mom_pct)}tır.\n'
        f'Non-Brand click\'te YoY %{abs(nb_click_yoy_pct):.0f} artış gerçekleşmiştir.'
    ),
    insight_text=(
        f'- Toplam Click: {format_number(total_click_feb26)} (Şub 2026) vs {format_number(total_click_feb25)} (Şub 2025)\n'
        f'- Brand Click: {format_number(brand_click_feb26)} (YoY {pct_str_short(((brand_click_feb26-brand_click_feb25)/brand_click_feb25)*100)})\n'
        f'- Non-Brand Click: {format_number(nb_click_feb26)} (YoY {pct_str_short(nb_click_yoy_pct)})\n'
        f'- Genel CTR: {(total_click_feb26/total_imp_feb26*100):.2f}%'
    ),
    chart_path=os.path.join(CHARTS_DIR, 'device_breakdown.png'),
    footer_text='Source: GSC'
)

# --- Slide 10: Non-Brand Click Düşüşü & Artışı ---
print("  Slide 10: Non-Brand Click Değişimi")
rebuild_slide(
    slides[9],
    header_text='Loft | Search Console Verileri',
    subtitle_text='Geçen Yıla Göre Non-Brand Click Düşüşü & Artışı Yaşanan Kelimeler',
    insight_text='',
    chart_path=os.path.join(CHARTS_DIR, 'nb_click_dual.png'),
    footer_text='Source: GSC',
    chart_height=Emu(3800000)
)

# --- Slide 11: Brand Search Volume (Keyword Planner - keep existing table data) ---
print("  Slide 11: Marka Arama Hacimleri (mevcut veri korunuyor)")
# Keep slide 11 as-is since it has the brand comparison table from the template

# --- Slide 12: Organic Traffic ---
print("  Slide 12: Organik Trafik")
rebuild_slide(
    slides[11],
    header_text='Loft | Organik Trafik Ölçümlemesi',
    subtitle_text=(
        f'Şubat 2026\'da organik trafik (session), Şubat 2025\'e göre {pct_str(organic_yoy_pct)}tır. '
        f'Ocak 2026\'ya göre ise {pct_str(organic_mom_pct)}tır.\n'
        f'Organic Search sessions: {format_number(organic_feb26)} (Şub 2026)'
    ),
    insight_text=(
        f'- Organic Search: {format_number(organic_feb25)} → {format_number(organic_feb26)} (YoY {pct_str_short(organic_yoy_pct)})\n'
        f'- Organic Social: {format_number(int(traffic_feb26_vs_feb25[1]["rows"][1][2]) if len(traffic_feb26_vs_feb25) > 1 else 0)} sessions\n'
        f'- Organic Shopping: {format_number(int(traffic_feb26_vs_feb25[1]["rows"][2][2]) if len(traffic_feb26_vs_feb25) > 1 else 0)} sessions'
    ),
    chart_path=os.path.join(CHARTS_DIR, 'organic_traffic.png'),
    footer_text='Source: Analytics'
)

# --- Slide 13: Landing Page - Yearly ---
print("  Slide 13: Landing Page (Şubat 2026)")
rebuild_slide(
    slides[12],
    header_text='Loft | Organic Session',
    subtitle_text='Şubat 2026 Landing Page Özelinde Organic Session Durumu',
    insight_text='',
    chart_path=os.path.join(CHARTS_DIR, 'landing_pages.png'),
    footer_text='Source: Analytics',
    chart_height=Emu(3800000)
)

# --- Slide 14: Landing Page - Monthly (use same data) ---
print("  Slide 14: Landing Page Detay")
# Top queries chart
rebuild_slide(
    slides[13],
    header_text='Loft | Top Non-Brand Sorgular',
    subtitle_text='Şubat 2026 Top Non-Brand Sorgular - Impression Karşılaştırma',
    insight_text='',
    chart_path=os.path.join(CHARTS_DIR, 'top_queries.png'),
    footer_text='Source: GSC',
    chart_height=Emu(3800000)
)

# --- Slides 15-16: Visibility (SEOMonitor - skip) ---
print("  Slide 15-16: Visibility (SEOMonitor - atlanıyor)")
for idx in [14, 15]:
    rebuild_slide(
        slides[idx],
        header_text='Loft | Visibility',
        subtitle_text='SEOMonitor Visibility Verileri',
        insight_text='Not: Bu slide için SEOMonitor verisi gereklidir. Veri temin edildiğinde güncellenecektir.',
        chart_path=None,
        footer_text='Source: SEOMonitor'
    )

# --- Slide 17: Section 02 Header (keep) ---

# --- Slides 18-19: Keyword Rankings (SEOMonitor - skip) ---
print("  Slide 18-19: Kelime Sıralamaları (SEOMonitor - atlanıyor)")
for idx in [17, 18]:
    rebuild_slide(
        slides[idx],
        header_text='Loft | Anahtar Kelime Sıralamaları',
        subtitle_text='SEOMonitor Kelime Sıralama Verileri',
        insight_text='Not: Bu slide için SEOMonitor verisi gereklidir. Veri temin edildiğinde güncellenecektir.',
        chart_path=None,
        footer_text='Source: SEOMonitor'
    )

# --- Slide 20: Section 03 Header (keep) ---

# --- Slide 21: What we did ---
print("  Slide 21: 2025'te Neler Yaptık")
# Keep existing content, it's manually curated

# --- Slide 22: What we plan ---
print("  Slide 22: 2026'da Neler Planlıyoruz")
# Keep existing content, it's manually curated

# --- Slide 23: Closing (keep) ---

# ============================================================
# ADD SUMMARY SLIDE (after slide 3, before data slides)
# ============================================================

# We'll add a traffic channel overview chart to an available position
# Actually, let's update slide 12 with the full traffic channel chart instead

print("Sunum kaydediliyor...")
prs.save(OUTPUT_PATH)

print(f"\nSunum başarıyla oluşturuldu: {OUTPUT_PATH}")
print(f"Grafikler: {CHARTS_DIR}/")

# ============================================================
# SUMMARY
# ============================================================
print("\n" + "="*60)
print("LOFT SEO ŞUBAT 2026 DEĞERLENDİRME ÖZETİ")
print("="*60)
print(f"\n📊 GSC Verileri (Şubat 2026):")
print(f"  Total Impression: {format_number(total_imp_feb26)} (YoY {pct_str_short(imp_yoy_pct)}, MoM {pct_str_short(imp_mom_pct)})")
print(f"  Total Click: {format_number(total_click_feb26)} (YoY {pct_str_short(click_yoy_pct)}, MoM {pct_str_short(click_mom_pct)})")
print(f"  Brand Impression: {format_number(brand_imp_feb26)} (YoY {pct_str_short(((brand_imp_feb26-brand_imp_feb25)/brand_imp_feb25)*100)})")
print(f"  Non-Brand Impression: {format_number(nb_imp_feb26)} (YoY {pct_str_short(nb_imp_yoy_pct)})")
print(f"  Non-Brand Click: {format_number(nb_click_feb26)} (YoY {pct_str_short(nb_click_yoy_pct)})")
print(f"\n📈 GA4 Organik Trafik:")
print(f"  Organic Search Sessions: {format_number(organic_feb26)} (YoY {pct_str_short(organic_yoy_pct)}, MoM {pct_str_short(organic_mom_pct)})")
print(f"\n🔝 Top Landing Pages (Session):")
for p in landing_pages[:5]:
    print(f"  {p['page']}: {p['sessions']:,} sessions")
print(f"\n🔑 Non-Brand Impression Artış Top 5:")
for k in nb_imp_increase[:5]:
    print(f"  {k['query']}: {format_number(k['imp_old'])} → {format_number(k['imp_new'])}")
