#!/usr/bin/env python3
"""
inbound_deck.py - Inbound Design System PPTX uretici.

deck.json (bildirimsel sahne tanimi) -> duzenlenebilir PPTX.

Koordinat sistemi: 1280x720 px sahne. 96 DPI'da tam olarak 13.333x7.5 inch,
yani standart PPTX 16:9. 1 px = 9525 EMU. Bu sayede Design System'in HTML
slaytlari ile PPTX birebir ayni izgarayi paylasir; build_html_preview.py ayni
deck.json'dan onizleme uretir ve iki cikti piksel piksel ortusur.

Font olcumu gercek TTF uzerinden PIL ile yapilir (assets/design-system/fonts).
Bu yuzden baslik auto-shrink, tablo kolon genisligi ve tasma tespiti tahmine
degil olcume dayanir - gecmis destelerdeki en sik teslim hatasi metin tasmasiydi.

Kullanim:
    python3 inbound_deck.py deck.json -o cikti.pptx
    python3 inbound_deck.py deck.json --check      # sadece dogrula, dosya yazma
"""

from __future__ import annotations

import argparse
import copy
import json
import os
import re
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ----------------------------------------------------------------------------
# Geometri
# ----------------------------------------------------------------------------

PX = 9525                 # 1 px @96dpi = 9525 EMU
STAGE_W, STAGE_H = 1280, 720
PT_PER_PX = 0.75          # px -> pt
PX_PER_PT = 4.0 / 3.0     # pt -> px


def px(v) -> Emu:
    return Emu(int(round(v * PX)))


# ----------------------------------------------------------------------------
# Design System token'lari (colors_and_type.css ile birebir)
# ----------------------------------------------------------------------------

C = {
    "coral":      "FF7B52",
    "coral_tint": "FFE3D8",
    "coral_deep": "E85F36",
    "teal":       "10332F",
    "teal_soft":  "1A4238",
    "mint":       "E8F5E9",
    "white":      "FFFFFF",
    "ink":        "10332F",
    "ink2":       "4A4A4A",
    "ink3":       "8A8A8A",
    "line":       "E0E0E0",
    "line_soft":  "F0EDE8",
    "red":        "D32F2F",
    "red_wash":   "FFCDD2",
    "green":      "2E7D32",
    "green_wash": "C8E6C9",
    "gold":       "F5A623",
    "gray_bar":   "4A4A4A",
    # VitrA destesinden olculmus deger token'lari
    "paper":      "FEFEF7",   # koyu/coral zemin uzerindeki metin (saf beyaz degil)
    "paper_bg":   "FEFFFA",   # ajanda slayti zemini
    "sep_num":    "254E49",   # ayrac numerali (teal uzerinde soluk)
}

F_DISPLAY = "Bricolage Grotesque"
F_BODY = "Outfit"

# Slayt izgarasi (Design System slides/*.html padding degerleri)
M_L, M_R = 60, 60                 # icerik slaytlari yan bosluk
BREADCRUMB_XY = (48, 28)
TITLE_TOP = 88
BODY_BOTTOM = 636                 # altinda logo/kaynak seridi
LOGO_XY = (44, 652)
LOGO_WH = 36
SOURCE_XY = (100, 658)

# ---- Kapak / ajanda / ayrac: VitrA destesinden olculmus degerler ----------
# Kaynak deste 2560x1440 olceginde; asagidaki degerler 1280x720 sahneye
# donusturulmus hali (punto ve koordinat /2).
COVER_TITLE_PT, COVER_TITLE_W = 47.5, 600   # Bricolage SemiBold
COVER_SUB_PT,   COVER_SUB_W   = 43.0, 600
COVER_TITLE_Y = 275
COVER_SUB_GAP = 22                 # baslik blogu ile donem satiri arasi
COVER_ART_W = 403                  # sol dekoratif big-O genisligi
COVER_WM = (563, 635, 153, 32)     # wordmark x,y,w,h

AGENDA_PANEL_W = 640               # sol coral panel
AGENDA_TITLE_PT, AGENDA_TITLE_W = 48.0, 400
AGENDA_TITLE_Y = 313
AGENDA_EYEBROW_XY = (27, 20)
AGENDA_LOGO = (31, 632, 52, 51)
AGENDA_LIST_X, AGENDA_LIST_W = 665, 573
AGENDA_ITEM_PT = 20.0
AGENDA_ITEM_LH = 0.95
AGENDA_ITEM_NUM_GAP = 8    # numara ile etiket arasi
AGENDA_ITEM_GAP = 26       # maddeler arasi

SEP_NUM_PT, SEP_NUM_W = 150.0, 800          # Bricolage ExtraBold (200pt'ten %25 kucultuldu)
SEP_NUM_CX = 146                            # numeral yatay merkezi (sabit)
SEP_TITLE_PT, SEP_TITLE_W = 37.0, 800
SEP_ACC_W, SEP_ACC_H = 43, 11               # accent cizgi (coral, VitrA olcusu)
SEP_ACC_GAP = 24                            # accent ile baslik blogu arasi
SEP_TITLE_MAX_W = 1080                      # bunu asan baslik alt satira kayar

# Punto olcegi. Govde metni okunabilirlik icin yukseltildi; 12pt tavan olarak
# sabit - bunun uzeri slaytta fazla buyuk duruyor. Bos alan kalacaksa yaziyi
# buyutmek tercih edilir, kucuk punto ile bosluk birakmak degil.
PT = {
    "cover":      66,
    "section":    45,
    "section_num": 157,
    "h1":         27,
    "h2":         21,
    "h3":         18,
    "h4":         13.5,     # blok basligi (baslik, govde degil)
    "lead":       12,
    "body":       12,       # insight, metin, not kutusu govdesi - TAVAN
    "sm":         11,       # panel maddeleri, ikincil govde
    "table":      10.5,     # tablo govdesi
    "xs":         10,       # tablo basligi, kpi delta
    "pill":       8.5,      # kaynak pill
    "micro":      9,        # dipnot, grafik etiketi, kpi etiketi
}
BODY_PT_MAX = 12            # govde metni bu puntoyu asmaz

# ----------------------------------------------------------------------------
# Font olcumu (gercek TTF, PIL)
# ----------------------------------------------------------------------------

_FONT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                         "..", "assets", "design-system", "fonts")
# Variable TTF: tek dosya tum agirliklari tasiyor. wght ekseni istenen agirliga
# cekilir; Bricolage'da ayrica opsz (optik boyut) ekseni var ve punto degerine
# gore 12-96 arasina kirpilir.
_VAR = {F_DISPLAY: "BricolageGrotesque-var.ttf", F_BODY: "Outfit-var.ttf"}
_cache: dict = {}
_HAVE_PIL = True
FONT_LOAD_ERROR = None          # yuklenemezse sebebi burada durur
try:
    from PIL import ImageFont
except Exception:          # pragma: no cover
    _HAVE_PIL = False
    FONT_LOAD_ERROR = "Pillow kurulu degil"


def _font(family: str, weight: int, size_px: float):
    """Variable fontu istenen agirlik ve optik boyutta dondurur."""
    global FONT_LOAD_ERROR
    key = (family, weight, round(size_px, 1))
    if key in _cache:
        return _cache[key]
    if not _HAVE_PIL:
        return None
    path = os.path.normpath(os.path.join(_FONT_DIR, _VAR.get(family, "")))
    f = None
    try:
        f = ImageFont.truetype(path, max(1, int(round(size_px))))
        axes = f.get_variation_axes()
        vals = []
        for a in axes:
            nm = a["name"]
            nm = nm.decode("utf-8", "ignore") if isinstance(nm, bytes) else str(nm)
            if "eight" in nm or nm == "wght":
                vals.append(max(a["minimum"], min(a["maximum"], weight)))
            elif "ptical" in nm or nm == "opsz":
                pt_ = size_px * PT_PER_PX
                vals.append(max(a["minimum"], min(a["maximum"], pt_)))
            else:
                vals.append(a["default"])
        if vals:
            f.set_variation_by_axes(vals)
    except Exception as e:
        if FONT_LOAD_ERROR is None:
            FONT_LOAD_ERROR = f"{os.path.basename(path)}: {e}"
        f = None
    _cache[key] = f
    return f


def text_w(s: str, pt: float, family: str = F_BODY, bold: bool = False,
           weight: int = None) -> float:
    """Metin genisligi (px). Font yuklenemezse karakter genisligi yaklasimina duser."""
    if not s:
        return 0.0
    size_px = pt * PX_PER_PT
    f = _font(family, weight if weight else (700 if bold else 400), size_px)
    if f is None:
        return len(s) * size_px * 0.52
    try:
        bbox = f.getbbox(s)
        return float(bbox[2] - bbox[0])
    except Exception:
        return len(s) * size_px * 0.52


# PIL'in advance genislikleri ile tarayicinin yerlesimi arasinda kucuk bir fark
# var (kerning ve letter-spacing kaynakli); tarayici bazen bir kelime once
# sariyor. Bu pay, olcumu muhafazakar tarafa cekerek PPTX ureticisinin
# tarayicidan ONCE uyarmasini saglar - iki cikti ayni sinirda okunur.
WRAP_SAFETY = 0.985


def wrap_lines(s: str, max_w: float, pt: float, family: str = F_BODY,
               bold: bool = False, safety: float = WRAP_SAFETY,
               weight: int = None) -> list:
    """Kelime bazli sarma; olculmus genisliklere gore satirlara boler."""
    max_w = max_w * safety
    words, lines, cur = s.split(), [], ""
    for w in words:
        trial = (cur + " " + w).strip()
        if cur and text_w(trial, pt, family, bold, weight) > max_w:
            lines.append(cur)
            cur = w
        else:
            cur = trial
    if cur:
        lines.append(cur)
    return lines or [""]


def fit_pt(s: str, max_w: float, start_pt: float, min_pt: float = 15,
           family: str = F_DISPLAY, bold: bool = True, weight: int = None) -> float:
    """Tek satira sigana kadar puntoyu kucultur (Design System: baslik sarmaz)."""
    p = start_pt
    while p > min_pt and text_w(s, p, family, bold, weight) > max_w:
        p -= 0.5
    return p


# ----------------------------------------------------------------------------
# Zengin metin: {b:..} bold, {g:..} yesil, {r:..} kirmizi, {c:..} coral, {n:..} soluk
# ----------------------------------------------------------------------------

_TAG = re.compile(r"\{([bgrcn]):([^{}]*)\}")
_STYLE = {
    "b": dict(bold=True, color="ink",  family=F_DISPLAY),
    "g": dict(bold=True, color="green", family=F_DISPLAY),
    "r": dict(bold=True, color="red",   family=F_DISPLAY),
    "c": dict(bold=True, color="coral", family=F_DISPLAY),
    "n": dict(bold=False, color="ink3", family=F_BODY),
}


def parse_runs(s: str, base_color: str = "ink", base_family: str = None) -> list:
    """'Organik {g:+%8 artmistir}.' -> [(metin, stil), ...]

    Isaretsiz metin ne font ailesi ne de kalinlik tasir; ikisini de cizildigi
    kutudan devralir. Onceden F_BODY yaziliyordu ve
    run duzeyi deger kutunun 'family' parametresini eziyordu - ayrac basliklari
    ile ajanda maddeleri F_DISPLAY istendigi halde Outfit basiliyordu
    (bkz. tuzaklar 3.6i). Etiketli parcalar (b/g/r/c/n) kendi ailelerini
    tasimaya devam eder.
    """
    out, pos = [], 0
    for m in _TAG.finditer(s or ""):
        if m.start() > pos:
            out.append((s[pos:m.start()],
                        dict(bold=None, color=base_color, family=base_family)))
        st = dict(_STYLE[m.group(1)])
        if st["color"] == "ink" and base_color != "ink":
            st["color"] = base_color
        out.append((m.group(2), st))
        pos = m.end()
    if pos < len(s or ""):
        out.append((s[pos:], dict(bold=None, color=base_color, family=base_family)))
    return out or [("", dict(bold=None, color=base_color, family=base_family))]


def plain(s: str) -> str:
    return _TAG.sub(r"\2", s or "")


# ----------------------------------------------------------------------------
# Alcak seviye cizim yardimcilari
# ----------------------------------------------------------------------------

def _no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, fill=None, radius=None, line=None, line_w=1):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, px(x), px(y), px(w), px(h))
    if radius:
        # adj = yaricap / kisa kenar
        try:
            s.adjustments[0] = min(0.5, radius / max(1.0, min(w, h)))
        except Exception:
            pass
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor.from_string(C.get(fill, fill))
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = RGBColor.from_string(C.get(line, line))
        s.line.width = Pt(line_w)
    else:
        _no_line(s)
    s.shadow.inherit = False
    if s.has_text_frame:
        s.text_frame.text = ""
    return s


def hline(slide, x, y, w, color="line", weight=1):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w),
                               Pt(weight))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor.from_string(C.get(color, color))
    _no_line(s)
    s.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h, runs, pt=PT["body"], family=F_BODY,
            color="ink", bold=False, align="l", anchor="t", line_pct=1.5,
            space_after=0, wrap=True):
    """runs: str | [(metin, stil)] | [ [(metin,stil)], ... ] (paragraf listesi)"""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    # python-pptx varsayilan olarak <a:spAutoFit/> yazar; kutu metne gore
    # buyur. Google Slides bu ayari kendi metrikleriyle yeniden uygular ve
    # kutular kayar/boyut degistirir. Geometri burada birebir hesaplandigi
    # icin otomatik boyutlama kapatilir (bkz. tuzaklar 3.6h).
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}[anchor]

    if isinstance(runs, str):
        paras = [parse_runs(runs, color, family)]
    elif runs and isinstance(runs[0], tuple):
        paras = [runs]
    else:
        paras = runs

    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                       "r": PP_ALIGN.RIGHT, "j": PP_ALIGN.JUSTIFY}[align]
        p.line_spacing = line_pct
        if space_after:
            p.space_after = Pt(space_after)
        for txt, st in (para or [("", {})]):
            r = p.add_run()
            r.text = txt
            fnt = r.font
            fnt.size = Pt(pt)
            fnt.name = st.get("family") or family
            _b = st.get("bold")
            fnt.bold = bool(bold if _b is None else _b)
            fnt.color.rgb = RGBColor.from_string(
                C.get(st.get("color", color), st.get("color", color)))
    return tb


def picture(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path):
        return None
    kw = {}
    if w:
        kw["width"] = px(w)
    if h:
        kw["height"] = px(h)
    return slide.shapes.add_picture(path, px(x), px(y), **kw)


def set_bg(slide, color):
    """Slayt zeminini doldurur (p:bg)."""
    xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr><a:solidFill><a:srgbClr val="%s"/></a:solidFill>'
        '<a:effectLst/></p:bgPr></p:bg>' % C.get(color, color)
    )
    from lxml import etree
    bg = etree.fromstring(xml)
    spTree = slide.shapes._spTree
    spTree.getparent().insert(0, bg)


# ----------------------------------------------------------------------------
# Slayt kromu: breadcrumb, baslik, alt baslik, logo, kaynak, dipnot
# ----------------------------------------------------------------------------

class Ctx:
    """Uretim baglami: varlik yollari, uyarilar, olcum kayitlari."""

    def __init__(self, assets_dir, base_dir):
        self.assets = assets_dir
        self.base = base_dir
        self.warnings = []
        self.boxes = []          # (slayt_no, etiket, x, y, w, h) - QA icin
        self.overflowed = set()  # (slayt_no, blok_tipi) - mukerrer uyari onleme

    def logo(self, name):
        return os.path.join(self.assets, "design-system", "logos", name)

    def slide_img(self, name):
        """Kapak / ajanda / ayrac gorselleri (VitrA destesinden cikarilmis)."""
        return os.path.join(self.assets, "design-system", "vitra-slides", name)

    def resolve(self, p):
        if not p:
            return p
        return p if os.path.isabs(p) else os.path.join(self.base, p)

    def warn(self, msg):
        self.warnings.append(msg)


def chrome(slide, spec, ctx, idx, dark=False):
    """Breadcrumb + baslik + alt baslik + logo + kaynak pill. Govde ust y dondurur."""
    inv = dark
    avail = STAGE_W - M_L - M_R

    if spec.get("breadcrumb"):
        bc = spec["breadcrumb"]
        parts = bc if isinstance(bc, list) else [bc]
        runs = []
        for i, part in enumerate(parts):
            if i:
                runs.append(("  |  ", dict(bold=False, color="coral", family=F_BODY)))
            runs.append((str(part), dict(bold=(i == 0), color="coral",
                                         family=F_DISPLAY if i == 0 else F_BODY)))
        textbox(slide, BREADCRUMB_XY[0], BREADCRUMB_XY[1], STAGE_W - 96, 18,
                runs, pt=PT["xs"], line_pct=1.0)

    y = TITLE_TOP
    if spec.get("title"):
        t = plain(spec["title"])
        p = fit_pt(t, avail, PT["h1"])
        if p < PT["h1"]:
            ctx.warn(f"S{idx}: baslik tek satira sigmasi icin {PT['h1']}pt -> {p}pt kucultuldu")
        textbox(slide, M_L, y, avail, p * PX_PER_PT * 1.15, t, pt=p,
                family=F_DISPLAY, bold=True,
                color="white" if inv else "ink", line_pct=1.05, wrap=False)
        y += p * PX_PER_PT * 1.15 + 10

    if spec.get("subtitle"):
        sub = spec["subtitle"]
        lines = wrap_lines(plain(sub), avail, PT["lead"])
        h = len(lines) * PT["lead"] * PX_PER_PT * 1.45
        textbox(slide, M_L, y, avail, h, sub, pt=PT["lead"],
                color="white" if inv else "ink2", line_pct=1.45)
        y += h + 16

    logo_name = "inbound-o-white.png" if inv else "inbound-o-teal.png"
    picture(slide, ctx.logo(logo_name), LOGO_XY[0], LOGO_XY[1], w=LOGO_WH, h=LOGO_WH)

    if spec.get("source"):
        src = spec["source"]
        if not src.strip().lower().startswith("kaynak"):
            src = "Kaynak: " + src.strip()
        w = text_w(src, PT["pill"], F_DISPLAY, True) + 24
        s = rect(slide, SOURCE_XY[0], SOURCE_XY[1], w, 22, fill="coral", radius=8)
        tf = s.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = src
        r.font.size = Pt(PT["pill"])
        r.font.name = F_DISPLAY
        r.font.bold = True
        r.font.color.rgb = RGBColor.from_string(C["white"])

    return y


def footnotes(slide, notes, ctx, idx):
    """Yildizli dipnotlar: govde altinda, kaynak seridinin ustunde."""
    if not notes:
        return BODY_BOTTOM
    notes = notes if isinstance(notes, list) else [notes]
    avail = STAGE_W - M_L - M_R
    paras, total = [], 0
    for n in notes:
        ls = wrap_lines(plain(n), avail, PT["micro"])
        total += len(ls) * PT["micro"] * PX_PER_PT * 1.45
        paras.append(parse_runs(n, "ink3"))
    y = BODY_BOTTOM - total
    textbox(slide, M_L, y, avail, total, paras, pt=PT["micro"],
            color="ink3", line_pct=1.45, space_after=2)
    return y - 12


# ----------------------------------------------------------------------------
# Blok: tablo
# ----------------------------------------------------------------------------

_NUMISH = re.compile(r"^[\s+\-]*[\d.,]+\s*(%|p|x|K|M|₺|pt)?[\s]*$|^%[\s]*[+\-]?[\d.,]+")


def _delta_kind(v: str):
    """'+%8.1' -> 'pos', '-%37' -> 'neg', digerleri None."""
    s = (v or "").strip()
    if not s or s in {"-", "—", "n/a"}:
        return None
    if s.startswith("+"):
        return "pos"
    if s.startswith("-") and re.search(r"\d", s):
        return "neg"
    return None


def _cell_num(v: str):
    """Hucre metninden sayi cikarir: '226.0K' -> 226000, '-%5' -> -5,
    '1.2M' -> 1200000, '54.1' -> 54.1. Sayi yoksa None."""
    s = plain(str(v or "")).strip().replace("%", "").replace(" ", "")
    if not s or s in {"-", "—", "n/a"}:
        return None
    mult = 1.0
    if s[-1:] in {"K", "k"}:
        mult, s = 1_000.0, s[:-1]
    elif s[-1:] in {"M", "m"}:
        mult, s = 1_000_000.0, s[:-1]
    elif s[-1:] in {"p", "x"}:
        s = s[:-1]
    s = s.replace(",", ".")
    m = re.fullmatch(r"[+-]?\d*\.?\d+", s)
    if not m:
        return None
    try:
        return float(s) * mult
    except ValueError:
        return None


def heat_cells(b, rows, ncol, dcols):
    """Satir bazli isi haritasi: her satirin en yuksek degeri yesil, en dusugu
    kirmizi. VitrA/Ozdilekteyim destelerindeki aylik hacim tablolarinin kullanimi.

    `heat: true` ile acilir. Kapsam varsayilani: etiket kolonu (0) ve delta
    kolonlari haric tum kolonlar; `heat_cols` ile acikca verilebilir.
    `heat_rows` verilmezse tum veri satirlari taranir.

    Iki renderer da bu fonksiyonu cagirir - renk karari tek yerde uretilir.
    En az uc sayisal deger yoksa satir atlanir (iki degerde "en yuksek/en dusuk"
    bilgi tasimaz, delta kolonu isini zaten yapar).
    """
    if not b.get("heat"):
        return {}
    cols = b.get("heat_cols")
    if cols is None:
        cols = [c for c in range(1, ncol) if c not in dcols]
    cols = [c for c in cols if 0 <= c < ncol]
    hrows = b.get("heat_rows")
    if hrows is None:
        hrows = range(len(rows))
    else:
        hrows = [r if r >= 0 else len(rows) + r for r in hrows]
    inv_rows = {r if r >= 0 else len(rows) + r
                for r in (b.get("heat_invert_rows") or [])}
    out = {}
    for ri in hrows:
        if not (0 <= ri < len(rows)):
            continue
        r = rows[ri]
        vals = [(ci, _cell_num(r[ci])) for ci in cols if ci < len(r)]
        vals = [(ci, v) for ci, v in vals if v is not None]
        if len(vals) < 3:
            continue
        hi = max(v for _, v in vals)
        lo = min(v for _, v in vals)
        if hi == lo:
            continue
        # Esit degerlerin tamami isaretlenir: 110.0K iki ayda da goruluyorsa
        # ikisi de en yuksek aydir, ilk gorulen secilmez. Ancak cok sayida
        # deger esitse (kovalanmis veri) isaretleme bilgi tasimaz - yarim satir
        # boyanmis olur. Bu durumda o uc isaretlenmez.
        ties = int(b.get("heat_max_ties", 2))
        n_hi = sum(1 for _, v in vals if v == hi)
        n_lo = sum(1 for _, v in vals if v == lo)
        # Ortalama pozisyon gibi kucuk degerin iyi oldugu satirlarda renk ters
        # cevrilir: en dusuk deger yesil olur.
        good_hi = ri not in inv_rows
        for ci, v in vals:
            if v == hi and n_hi <= ties:
                out[(ri, ci)] = "pos" if good_hi else "neg"
            elif v == lo and n_lo <= ties:
                out[(ri, ci)] = "neg" if good_hi else "pos"
    return out


def table_layout(b, w):
    """
    Tablo geometrisini hesaplar: kolon genislikleri, satir yukseklikleri, baslik
    yuksekligi. HTML onizleme de bu fonksiyonu kullanir - aksi halde tarayici
    satir yuksekligini padding+icerikten turetiyor ve tablolar PPTX'ten daha
    uzun render ediliyordu (9 satirda ~45px fark). Tek kaynak, iki cikti.
    """
    head = list(b.get("head") or [])
    rows = b.get("rows") or []
    ncol = max([len(head)] + [len(r) for r in rows] or [1])
    head += [""] * (ncol - len(head))
    pad = 12
    th_pt, td_pt = PT["xs"], b.get("font_pt", PT["table"])

    if b.get("col_w"):
        tot = float(sum(b["col_w"]))
        widths = [w * c / tot for c in b["col_w"]]
    else:
        # Her hucre, cizilecegi font ve agirlikla olculur. Delta hucreleri ve
        # kalin satirlar F_DISPLAY bold ile render ediliyor; olcumu F_BODY ile
        # yapmak kolonu daraltip "+%92.8" gibi degerlerin iki satira kirilmasina
        # yol aciyordu.
        bold_rows_m = {r if r >= 0 else len(rows) + r
                       for r in (b.get("bold_rows") or [])}
        delta_m = set(b.get("delta_cols") or [])
        if not delta_m:
            for ci in range(ncol):
                vals = [r[ci] for r in rows if ci < len(r)]
                if vals and sum(1 for v in vals if _delta_kind(str(v))) >= max(1, len(vals) // 2):
                    delta_m.add(ci)
        need = []
        for ci in range(ncol):
            mx = text_w(plain(str(head[ci])), th_pt, F_DISPLAY, True)
            for ri, r in enumerate(rows):
                if ci >= len(r):
                    continue
                val = plain(str(r[ci]))
                strong = (ci in delta_m and _delta_kind(val)) or ri in bold_rows_m
                mx = max(mx, text_w(val, td_pt,
                                    F_DISPLAY if strong else F_BODY, bool(strong)))
            need.append(mx + pad * 2)
        # Ilk kolon (etiket/metrik) artan genisligi tek basina yutmamali: "Yil"
        # gibi kisa bir baslik tablonun yarisini kaplayabiliyordu. Etiket kolonu
        # kendi ihtiyacini alir, tavani tablo genisliginin first_col_max'i kadar;
        # kalan bosluk sayisal kolonlara esit dagitilir.
        fixed = sum(need[1:])
        cap = w * float(b.get("first_col_max", 0.34))
        first = need[0] if need[0] <= cap else max(cap, 90.0)
        slack = w - first - fixed
        if slack > 0 and ncol > 1:
            widths = [first] + [n + slack / (ncol - 1) for n in need[1:]]
        else:
            widths = [first] + need[1:]
        if sum(widths) > w:
            # Once ilk kolon (etiket) daraltilir; etiket sarabilir, sayisal
            # kolonlar sarmaz. Yetmezse hepsi oransal daraltilir ve uyarilir.
            over = sum(widths) - w
            take = min(over, max(0.0, widths[0] - 90))
            widths[0] -= take
            if sum(widths) > w + 0.5:
                sc = w / sum(widths)
                widths = [wd * sc for wd in widths]

    row_h = b.get("row_h", 26)
    head_h = b.get("head_h", 30)
    row_hs = []
    for r in rows:
        n = len(wrap_lines(plain(str(r[0] if r else "")), widths[0] - pad * 2,
                           td_pt, safety=1.0)) if ncol else 1
        row_hs.append(max(row_h, n * td_pt * PX_PER_PT * 1.35 + 10))

    title_h = (PT["h4"] * PX_PER_PT * 1.3 + 8) if b.get("title") else 0
    return dict(head=head, rows=rows, ncol=ncol, pad=pad, th_pt=th_pt,
                td_pt=td_pt, widths=widths, row_hs=row_hs, head_h=head_h,
                title_h=title_h, total_h=title_h + head_h + sum(row_hs))


def block_table(slide, b, x, y, w, ctx, idx):
    """
    b = {type:"table", head:[...], rows:[[...]], align:"lccc", delta_cols:[3],
         wash:true, highlight_rows:[2], bold_rows:[-1], col_w:[40,20,20,20],
         row_h:26, head_h:30, font_pt:9.5, title:"opsiyonel"}
    Geometri table_layout ile hesaplanir; HTML onizleme ayni fonksiyonu kullanir.
    """
    L = table_layout(b, w)
    head, rows, ncol = L["head"], L["rows"], L["ncol"]
    pad, th_pt, td_pt = L["pad"], L["th_pt"], L["td_pt"]
    widths, row_hs, head_h = L["widths"], L["row_hs"], L["head_h"]

    align = b.get("align") or ("l" + "c" * (ncol - 1))
    align = (align + "c" * ncol)[:ncol]
    wash = b.get("wash", True)
    hl_rows = set(b.get("highlight_rows") or [])
    bold_rows = {r if r >= 0 else len(rows) + r for r in (b.get("bold_rows") or [])}

    delta_cols = set(b.get("delta_cols") or [])
    if not delta_cols:
        for ci in range(ncol):
            vals = [r[ci] for r in rows if ci < len(r)]
            if vals and sum(1 for v in vals if _delta_kind(str(v))) >= max(1, len(vals) // 2):
                delta_cols.add(ci)
    heat = heat_cells(b, rows, ncol, delta_cols)

    y0 = y
    if b.get("title"):
        textbox(slide, x, y, w, PT["h4"] * PX_PER_PT * 1.3, b["title"],
                pt=PT["h4"], family=F_DISPLAY, bold=True, color="ink", line_pct=1.2)
        y += PT["h4"] * PX_PER_PT * 1.3 + 8

    # baslik satiri (teal zemin + beyaz kalin) - Bolum 15.3 deck standardi
    rect(slide, x, y, w, head_h, fill="teal")
    cx = x
    for ci in range(ncol):
        textbox(slide, cx + pad, y, widths[ci] - pad * 2, head_h,
                plain(str(head[ci])), pt=th_pt, family=F_DISPLAY, bold=True,
                color="white", align=align[ci], anchor="m", line_pct=1.0)
        cx += widths[ci]
    y += head_h

    for ri, r in enumerate(rows):
        rh = row_hs[ri]
        if ri in hl_rows:
            rect(slide, x, y, w, rh, fill="coral_tint")
        cx = x
        for ci in range(ncol):
            val = str(r[ci]) if ci < len(r) else ""
            kind = _delta_kind(plain(val)) if ci in delta_cols else None
            if kind is None:
                kind = heat.get((ri, ci))
            col, bold = "ink", (ri in bold_rows)
            if kind == "pos":
                col, bold = "green", True
                if wash:
                    rect(slide, cx, y, widths[ci], rh, fill="green_wash")
            elif kind == "neg":
                col, bold = "red", True
                if wash:
                    rect(slide, cx, y, widths[ci], rh, fill="red_wash")
            textbox(slide, cx + pad, y, widths[ci] - pad * 2, rh, val,
                    pt=td_pt, family=F_DISPLAY if bold else F_BODY, bold=bold,
                    color=col, align=align[ci], anchor="m", line_pct=1.3,
                    wrap=(ci == 0))
            cx += widths[ci]
        hline(slide, x, y, w, "line", 0.75)
        y += rh

    hline(slide, x, y, w, "line", 0.75)
    ctx.boxes.append((idx, "table", x, y0, w, y - y0))
    return y - y0


# ----------------------------------------------------------------------------
# Blok: insight listesi
# ----------------------------------------------------------------------------

def block_insights(slide, b, x, y, w, ctx, idx):
    """b = {type:"insights", title?:"TREND OKUMALARI", items:[".."], dark?:false}"""
    items = b.get("items") or []
    dark = b.get("dark", False)
    base = "paper" if dark else "ink"
    y0, ind, gap = y, 26, b.get("gap", 12)

    if b.get("title"):
        textbox(slide, x, y, w, 20, plain(b["title"]).upper(), pt=PT["xs"],
                family=F_DISPLAY, bold=True, color="coral", line_pct=1.0)
        y += 24

    pt_ = min(b.get("font_pt", PT["body"]), BODY_PT_MAX)
    for it in items:
        lines = wrap_lines(plain(it), w - ind, pt_)
        h = len(lines) * pt_ * PX_PER_PT * 1.5
        # ok her zaman coral: Design System'de insight isareti marka turuncusu
        textbox(slide, x, y, 18, pt_ * PX_PER_PT * 1.5, "➔", pt=pt_,
                family=F_BODY, color="coral", line_pct=1.5)
        textbox(slide, x + ind, y, w - ind, h, it, pt=pt_, color=base, line_pct=1.5)
        y += h + gap

    ctx.boxes.append((idx, "insights", x, y0, w, y - y0))
    return y - y0


# ----------------------------------------------------------------------------
# Blok: KPI kartlari
# ----------------------------------------------------------------------------

def block_kpi(slide, b, x, y, w, ctx, idx):
    """
    b = {type:"kpi", cards:[{value:"595.4K", unit?:"+", label:"CLICK",
                             delta?:"+%12.4 YoY", accent?:"teal|coral"}],
         cols?:4, h?:132}
    """
    cards = b.get("cards") or []
    cols = b.get("cols") or min(4, max(1, len(cards)))
    gap = b.get("gap", 18)
    cw = (w - gap * (cols - 1)) / cols
    ch = b.get("h", 132)
    rows = (len(cards) + cols - 1) // cols
    _tasma_uyarildi = False

    for i, c in enumerate(cards):
        cx = x + (i % cols) * (cw + gap)
        cy = y + (i // cols) * (ch + gap)
        accent = c.get("accent") or ("coral" if (i == len(cards) - 1
                                                and len(cards) > 2) else "teal")
        rect(slide, cx, cy, cw, ch, fill=accent, radius=16)

        val, unit = str(c.get("value", "")), str(c.get("unit", "") or "")
        deltas = c.get("deltas")
        if not deltas and c.get("delta"):
            deltas = [{"label": "", "value": c["delta"]}]

        # Kart icerigi kart kutusunun DISINA tasmamali. Beyaz metin kartin
        # altina tastiginda beyaz zeminde gorunmez olur; PowerPoint'te fark
        # edilmeyip Google Slides'a aktarildiginda satir tamamen kaybolur
        # (bkz. tuzaklar 3.6g). Once genislige, sonra yuksekliğe gore kucultulur.
        vpt = c.get("pt") or 40
        while vpt > 20 and text_w(val + unit, vpt, F_DISPLAY, True) > cw - 32:
            vpt -= 1
        LBL_H, DLT_H, UST, ARA, ALT = 18, 18, 22, 8, 12

        def _icerik_h(p):
            return UST + p * PX_PER_PT * 1.05 + ARA + LBL_H + (2 + DLT_H if deltas else 0) + ALT

        while vpt > 20 and _icerik_h(vpt) > ch:
            vpt -= 1
        if _icerik_h(vpt) > ch and not _tasma_uyarildi:
            _tasma_uyarildi = True
            ctx.warn(f"KPI S{idx}: kart yuksekligi {ch:.0f}px, icerik "
                     f"{_icerik_h(vpt):.0f}px istiyor - 'h' degeri buyutulmeli; "
                     f"aksi halde delta satiri kartin disina taser ve beyaz "
                     f"zeminde gorunmez olur")
        # icerik blogu kart icinde dikeyde ortalanir
        vy = cy + max(UST * 0.5, (ch - (_icerik_h(vpt) - UST - ALT)) / 2)

        runs = [(val, dict(bold=True, color="white", family=F_DISPLAY))]
        textbox(slide, cx, vy, cw, vpt * PX_PER_PT * 1.05, runs, pt=vpt,
                family=F_DISPLAY, bold=True, color="white", align="c",
                line_pct=1.0, wrap=False)
        if unit:
            textbox(slide, cx, vy + vpt * PX_PER_PT * 0.12, cw,
                    vpt * PX_PER_PT * 0.6, unit, pt=vpt * 0.5, family=F_DISPLAY,
                    bold=True, color="white", align="c", line_pct=1.0, wrap=False)

        ly = vy + vpt * PX_PER_PT * 1.05 + ARA
        textbox(slide, cx + 12, ly, cw - 24, LBL_H,
                plain(str(c.get("label", ""))).upper(),
                pt=PT["micro"], color="white", align="c", line_pct=1.1)

        # Delta satiri: etiket once, deger sonra ("MoM  +%8.3    YoY  -%18.9").
        # Etiket normal agirlikta ve hafif soluk, deger kalin - okuma sirasi
        # "hangi karsilastirma" -> "ne kadar" seklinde olur.
        if deltas:
            runs = []
            for di, dd in enumerate(deltas):
                if di:
                    runs.append(("     ", dict(bold=False, color="white",
                                               family=F_BODY)))
                lb = str(dd.get("label", "")).strip()
                if lb:
                    runs.append((lb + "  ", dict(bold=False, color="white",
                                                 family=F_BODY)))
                runs.append((str(dd.get("value", "")),
                             dict(bold=True, color="white", family=F_DISPLAY)))
            textbox(slide, cx + 8, ly + 20, cw - 16, 18, runs, pt=PT["xs"],
                    family=F_BODY, color="white", align="c", line_pct=1.1,
                    wrap=False)

    h = rows * ch + (rows - 1) * gap
    ctx.boxes.append((idx, "kpi", x, y, w, h))
    return h


# ----------------------------------------------------------------------------
# Blok: bar / stacked bar (duzenlenebilir vektor sekiller)
# ----------------------------------------------------------------------------

def _fmt_axis(v):
    """bar blogu icin bicimleyici - tek kaynak _fmt_val (asagida tanimli).
    Ayri bir kopya tutulmaz; ucu de ayrisirsa ayni seride "665K" ile "711.7K"
    karisimi olusuyor (bkz. tuzaklar 3.6)."""
    return _fmt_val(v, "auto")


def block_bar(slide, b, x, y, w, ctx, idx):
    """
    b = {type:"bar", cats:["Oca","Sub",...],
         series:[{name:"2025", data:[...], color:"gold"},
                 {name:"2026", data:[...], color:"gray_bar"}],
         h?:280, value_labels?:true, stacked?:false, fmt?:"K"}
    Native pptx chart yerine sekil cizilir: PowerPoint kategori eksenini
    "1,2,3" gosterdigi icin (gecmis projelerde yasandi) ve tasarim sistemi
    gorunumu sekillerle birebir tutuluyor. Sekiller sunum icinden duzenlenebilir.
    """
    cats = b.get("cats") or []
    series = b.get("series") or []
    if not cats or not series:
        return 0
    stacked = b.get("stacked", False)
    h = b.get("h", 260)
    y0 = h_used = y

    if b.get("title"):
        textbox(slide, x, y, w, PT["h4"] * PX_PER_PT * 1.3, b["title"],
                pt=PT["h4"], family=F_DISPLAY, bold=True, line_pct=1.2)
        y += PT["h4"] * PX_PER_PT * 1.3 + 6

    # legend
    if len(series) > 1 or b.get("legend", True):
        lx = x
        for s in series:
            rect(slide, lx, y + 3, 10, 10, fill=s.get("color", "gray_bar"), radius=3)
            nm = str(s.get("name", ""))
            textbox(slide, lx + 15, y, 200, 14, nm, pt=PT["micro"], color="ink2",
                    line_pct=1.0, wrap=False)
            lx += 15 + text_w(nm, PT["micro"]) + 26
        y += CB_BAR_LEGEND_H

    lab_h = CB_CAT_H
    val_h = CB_VAL_H if b.get("value_labels", True) else 0
    plot_top = y + val_h
    plot_bot = y0 + h - lab_h
    plot_h = max(40, plot_bot - plot_top)
    # Deger etiketi bandi ayrilmis olsa bile plot alani cok kisaldiginda en
    # yuksek barin etiketi banda sigmaz. Uyarilir; slaytta h buyutulur.
    if val_h and plot_h < 90:
        ctx.warn(f"GRAFIK S{idx}: 'bar' blogunda plot yuksekligi {plot_h:.0f}px - "
                 f"deger etiketleri sikisiyor, h en az {90 + val_h + lab_h + CB_BAR_LEGEND_H:.0f} olmali")

    if stacked:
        tops = [sum(float(s["data"][i] or 0) for s in series) for i in range(len(cats))]
        vmax = max(tops) if tops else 1
    else:
        vmax = max((max(float(v or 0) for v in s["data"]) for s in series), default=1)
    vmax = vmax * 1.12 or 1
    # invert: kucuk degerin iyi oldugu metrikler (ortalama pozisyon). Bar
    # yukseklikleri ters cevrilir, boylece iyilesme yukari dogru okunur.
    inv = b.get("invert", False)

    # yatay izgara
    for k in range(1, 4):
        gy = plot_bot - plot_h * k / 4
        hline(slide, x, gy, w, "line_soft", 0.75)
    hline(slide, x, plot_bot, w, "line", 0.75)

    n = len(cats)
    slot = w / n
    if stacked:
        bw = min(b.get("bar_w", 46), slot * 0.62)
    else:
        bw = min(b.get("bar_w", 22), (slot * 0.72) / len(series))
    grp_w = bw if stacked else bw * len(series) + 4 * (len(series) - 1)

    for i, cat in enumerate(cats):
        cx = x + slot * i + (slot - grp_w) / 2
        if stacked:
            acc = 0.0
            for s in series:
                v = float(s["data"][i] or 0)
                bh = plot_h * v / vmax
                if bh >= 1:
                    rect(slide, cx, plot_bot - acc - bh, bw, bh,
                         fill=s.get("color", "gray_bar"))
                acc += bh
            if b.get("value_labels", True):
                tot = sum(float(s["data"][i] or 0) for s in series)
                textbox(slide, cx - 12, plot_bot - acc - 15, bw + 24, 14,
                        _fmt_axis(tot), pt=PT["micro"], family=F_DISPLAY,
                        bold=True, color="ink", align="c", line_pct=1.0, wrap=False)
        else:
            for si, s in enumerate(series):
                v = float(s["data"][i] or 0)
                bh = plot_h * (1 - v / vmax) if inv else plot_h * v / vmax
                bx = cx + si * (bw + 4)
                if bh >= 1:
                    rect(slide, bx, plot_bot - bh, bw, bh,
                         fill=s.get("color", "gray_bar"))
                if b.get("value_labels", True):
                    textbox(slide, bx - 14, plot_bot - bh - 15, bw + 28, 14,
                            _fmt_axis(v), pt=PT["micro"], family=F_DISPLAY,
                            bold=True, color="ink", align="c", line_pct=1.0,
                            wrap=False)
        textbox(slide, x + slot * i, plot_bot + 5, slot, 16, str(cat),
                pt=PT["micro"], color="ink2", align="c", line_pct=1.0, wrap=False)

    # Yatay cakisma: kategori etiketleri ve deger etiketleri wrap=False ile
    # ciziliyor, yani slot'a sigmazsa komsusunun uzerine tasar. Olcup uyaririz.
    _label_fit_check(ctx, idx, "bar", slot, [str(c) for c in cats],
                     [_fmt_axis(float(v or 0)) for s in series for v in s["data"]]
                     if b.get("value_labels", True) else [],
                     n_series=1 if stacked else len(series))

    ctx.boxes.append((idx, "bar", x, y0, w, h))
    return h


def _label_fit_check(ctx, idx, kind, slot, cats, vals, n_series=1):
    """Grafik etiketlerinin yatayda komsusuna binip binmedigini olcer.

    Etiketler wrap=False cizildigi icin slot'tan genis bir etiket sessizce
    yanindakinin uzerine tasar - PPTX'te de, onizlemede de. Taban kural:
    kategori etiketi slot'a, deger etiketi ise seri basina dusen paya sigmali.
    """
    pad = 4
    if cats:
        wide = max(cats, key=lambda t: text_w(t, PT["micro"]))
        need = text_w(wide, PT["micro"]) + pad
        if need > slot:
            ctx.warn(f"CAKISMA S{idx}: '{kind}' kategori etiketi '{wide}' "
                     f"{need:.0f}px yer istiyor, slot {slot:.0f}px - komsu etiketle "
                     f"cakisiyor; kategori sayisi azaltilmali veya etiket kisaltilmali")
    if vals:
        wide = max(vals, key=lambda t: text_w(t, PT["micro"]))
        need = text_w(wide, PT["micro"]) + pad
        avail = slot / max(1, n_series)
        if need > avail:
            ctx.warn(f"CAKISMA S{idx}: '{kind}' deger etiketi '{wide}' "
                     f"{need:.0f}px yer istiyor, seri basina {avail:.0f}px var - "
                     f"deger etiketleri ust uste biniyor; value_labels kapatilmali "
                     f"veya kategori sayisi azaltilmali")


def block_line(slide, b, x, y, w, ctx, idx):
    """
    b = {type:"line", cats:[...], series:[{name, data, color}], h?:260}
    Freeform polyline + nokta isaretleri; duzenlenebilir kalir.
    """
    cats = b.get("cats") or []
    series = b.get("series") or []
    if not cats or not series:
        return 0
    h = b.get("h", 260)
    y0 = y

    if b.get("title"):
        textbox(slide, x, y, w, PT["h4"] * PX_PER_PT * 1.3, b["title"],
                pt=PT["h4"], family=F_DISPLAY, bold=True, line_pct=1.2)
        y += PT["h4"] * PX_PER_PT * 1.3 + 6

    lx = x
    for s in series:
        rect(slide, lx, y + 4, 14, 3, fill=s.get("color", "coral"), radius=2)
        nm = str(s.get("name", ""))
        textbox(slide, lx + 19, y, 200, 14, nm, pt=PT["micro"], color="ink2",
                line_pct=1.0, wrap=False)
        lx += 19 + text_w(nm, PT["micro"]) + 26
    y += 20

    plot_top, plot_bot = y + 14, y0 + h - 20
    plot_h = max(40, plot_bot - plot_top)
    inv = b.get("invert", False)
    vals = [float(v or 0) for s in series for v in s["data"]]
    vmin, vmax = (min(vals), max(vals)) if vals else (0, 1)
    if vmax == vmin:
        vmax = vmin + 1
    span = (vmax - vmin) * 1.15
    base = vmin - (vmax - vmin) * 0.075

    for k in range(1, 4):
        hline(slide, x, plot_bot - plot_h * k / 4, w, "line_soft", 0.75)
    hline(slide, x, plot_bot, w, "line", 0.75)

    n = len(cats)
    step = w / max(1, n - 1) if n > 1 else w

    for s in series:
        col = C.get(s.get("color", "coral"), s.get("color", "coral"))
        pts = []
        for i, v in enumerate(s["data"][:n]):
            vx = x + step * i if n > 1 else x + w / 2
            frac = (float(v or 0) - base) / span
            vy = plot_bot - plot_h * ((1 - frac) if inv else frac)
            pts.append((vx, vy))
        if len(pts) > 1:
            ff = slide.shapes.build_freeform(px(pts[0][0]), px(pts[0][1]))
            ff.add_line_segments([(px(a), px(bb)) for a, bb in pts[1:]], close=False)
            shp = ff.convert_to_shape()
            shp.fill.background()
            shp.line.color.rgb = RGBColor.from_string(col)
            shp.line.width = Pt(2.25)
            shp.shadow.inherit = False
        for (vx, vy) in pts:
            d = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(vx - 3.5), px(vy - 3.5),
                                       px(7), px(7))
            d.fill.solid()
            d.fill.fore_color.rgb = RGBColor.from_string(col)
            _no_line(d)
            d.shadow.inherit = False

    for i, cat in enumerate(cats):
        vx = x + step * i if n > 1 else x + w / 2
        textbox(slide, vx - step / 2, plot_bot + 5, max(step, 40), 16, str(cat),
                pt=PT["micro"], color="ink2", align="c", line_pct=1.0, wrap=False)

    ctx.boxes.append((idx, "line", x, y0, w, h))
    return h


# ----------------------------------------------------------------------------
# Blok: combo (bar + cizgi, cift eksen) - Ozdilekteyim aylik metrik grafigi
# ----------------------------------------------------------------------------

CB_GUTTER = 54          # eksen etiketi icin sol/sag bosluk
CB_LEGEND_H = 22
CB_CAT_H = 20
CB_VAL_H = 16           # bar ustundeki deger etiketi bandi
CB_BAR_LEGEND_H = 20    # bar bloğundaki legend satiri


def _fmt_val(v, fmt):
    """Eksen ve etiket bicimleyici."""
    if fmt == "pct":
        return f"%{v:.2f}" if abs(v) < 10 else f"%{v:.1f}"
    if fmt == "pos":
        return f"{v:.1f}"
    a = abs(v)
    # K/M her zaman tek ondalikli: ayni seride "665K" ile "711.7K" karisimi olmaz
    # (Icerik Dili Rehberi 6.1 - tek rakam formati).
    if fmt == "M" or (fmt == "auto" and a >= 1_000_000):
        return f"{v/1_000_000:.1f}M"
    if fmt == "K" or (fmt == "auto" and a >= 1_000):
        return f"{v/1_000:.1f}K"
    if a and a < 10:
        return f"{v:.1f}"
    return f"{v:.0f}"


def _axis_scale(vals, invert, pad=1.15):
    """(vmin, vmax) - bar/cizgi eksen araligi. invert'te kucuk deger iyi."""
    if not vals:
        return 0.0, 1.0
    lo, hi = min(vals), max(vals)
    if invert:
        # pozisyon: en iyi (kucuk) deger ustte, 1'den baslar
        return max(0.5, lo * 0.85), hi * 1.12
    return 0.0, (hi * pad or 1.0)


def block_combo(slide, b, x, y, w, ctx, idx):
    """
    b = {type:"combo", h:300, cats:[...],
         series:[{kind:"bar", name, data, color, axis:"right", labels:"inside",
                  fmt:"M"},
                 {kind:"line", name, data, color, axis:"left", labels:"above",
                  fmt:"K", invert:false}],
         axis_labels: true}

    Ozdilekteyim destesindeki aylik metrik grafigi: bir metrik bar (deger
    etiketi barin icinde), diger metrik cizgi (deger etiketi noktanin ustunde),
    iki ayri y ekseni. Olcekleri farkli metrikleri tek grafikte okunur kilar.
    Pozisyon gibi kucuk degerin iyi oldugu seride invert:true - cizgi yukari
    ciktiginda iyilesme okunur, eksen etiketleri gercek degerleri gosterir.
    """
    cats = b.get("cats") or []
    series = b.get("series") or []
    if not cats or not series:
        return 0
    h = b.get("h", 300)
    y0 = y

    if b.get("title"):
        th = PT["h4"] * PX_PER_PT * 1.3
        textbox(slide, x, y, w, th, b["title"], pt=PT["h4"], family=F_DISPLAY,
                bold=True, line_pct=1.2)
        y += th + 6

    # legend - ortalanmis
    if b.get("legend", True):
        chips = []
        for s_ in series:
            nm = str(s_.get("name", ""))
            chips.append((s_, nm, text_w(nm, PT["micro"], F_BODY)))
        total = sum(18 + cw + 24 for _s, _n, cw in chips) - 24
        lx = x + (w - total) / 2
        for s_, nm, cw in chips:
            col = s_.get("color", "gray_bar")
            if s_.get("kind") == "line":
                rect(slide, lx, y + 8, 14, 3, fill=col, radius=2)
                d = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(lx + 4.5),
                                           px(y + 5.5), px(8), px(8))
                d.fill.solid()
                d.fill.fore_color.rgb = RGBColor.from_string(C.get(col, col))
                _no_line(d)
                d.shadow.inherit = False
            else:
                rect(slide, lx, y + 4, 12, 12, fill=col, radius=3)
            textbox(slide, lx + 18, y + 2, cw + 6, 14, nm, pt=PT["micro"],
                    color="ink2", line_pct=1.0, wrap=False)
            lx += 18 + cw + 24
        y += CB_LEGEND_H

    plot_top = y
    plot_bot = y0 + h - CB_CAT_H
    plot_h = max(50, plot_bot - plot_top)
    gut = CB_GUTTER if b.get("axis_labels", True) else 8
    px0, px1 = x + gut, x + w - gut
    pw = max(40, px1 - px0)
    n = len(cats)
    slot = pw / n

    # Seri-eksen baglantisi. Gecersiz bir axis degeri (ornegin "l"/"r") sessizce
    # elenir ve grafik izgarayla birlikte bos cizilir - tasma denetimine de
    # takilmaz. Bu yuzden once dogrulanir.
    for s_ in series:
        a = s_.get("axis", "left")
        if a not in ("left", "right"):
            ctx.warn(f"GRAFIK S{idx}: 'combo' serisi '{s_.get('name', '?')}' "
                     f"gecersiz axis degeri tasiyor: '{a}' - 'left' veya 'right' "
                     f"olmali. Seri cizilmeyecek.")
        if not [v for v in (s_.get("data") or []) if v not in (None, "")]:
            ctx.warn(f"GRAFIK S{idx}: 'combo' serisi '{s_.get('name', '?')}' "
                     f"veri tasimiyor - grafik bos cizilecek.")

    # eksen araliklari
    ax = {}
    for side in ("left", "right"):
        vals, inv, fmt = [], False, "auto"
        for s_ in series:
            if s_.get("axis", "left") != side:
                continue
            vals += [float(v) for v in s_.get("data", []) if v is not None]
            inv = inv or bool(s_.get("invert"))
            fmt = s_.get("fmt", fmt)
        if vals:
            lo, hi = _axis_scale(vals, inv)
            ax[side] = dict(lo=lo, hi=hi, inv=inv, fmt=fmt)

    if not ax:
        ctx.warn(f"GRAFIK S{idx}: 'combo' blogunda hicbir seri eksene baglanmadi - "
                 f"grafik yalnizca izgara olarak cizilecek. axis degerleri kontrol edilmeli.")

    def ypos(side, v):
        a = ax[side]
        frac = (float(v) - a["lo"]) / max(1e-9, a["hi"] - a["lo"])
        return plot_bot - plot_h * ((1 - frac) if a["inv"] else frac)

    # izgara + eksen etiketleri
    TICKS = 4
    for k in range(TICKS + 1):
        gy = plot_bot - plot_h * k / TICKS
        hline(slide, px0, gy, pw, "line_soft" if k else "line", 0.75)
        if not b.get("axis_labels", True):
            continue
        for side in ax:
            a = ax[side]
            frac = k / TICKS
            v = a["lo"] + (a["hi"] - a["lo"]) * ((1 - frac) if a["inv"] else frac)
            lbl = _fmt_val(v, a["fmt"])
            if side == "left":
                textbox(slide, x, gy - 7, gut - 8, 14, lbl, pt=PT["micro"],
                        color="ink3", align="r", line_pct=1.0, wrap=False)
            else:
                textbox(slide, px1 + 8, gy - 7, gut - 8, 14, lbl, pt=PT["micro"],
                        color="ink3", align="l", line_pct=1.0, wrap=False)

    # barlar
    for s_ in series:
        if s_.get("kind") != "bar":
            continue
        side = s_.get("axis", "left")
        if side not in ax:
            continue
        bw = min(b.get("bar_w", 40), slot * 0.62)
        col = s_.get("color", "gray_bar")
        lbls = s_.get("labels_text")
        for i, v in enumerate(s_["data"][:n]):
            v = float(v or 0)
            top = ypos(side, v)
            bh = max(1.0, plot_bot - top)
            bx = px0 + slot * i + (slot - bw) / 2
            rect(slide, bx, top, bw, bh, fill=col)
            if s_.get("labels") == "inside" and bh > 24:
                txt = lbls[i] if lbls and i < len(lbls) else _fmt_val(v, ax[side]["fmt"])
                textbox(slide, bx - 8, plot_bot - 20, bw + 16, 14, txt,
                        pt=PT["micro"], color="white", align="c", line_pct=1.0,
                        wrap=False)
            elif s_.get("labels") == "above":
                txt = lbls[i] if lbls and i < len(lbls) else _fmt_val(v, ax[side]["fmt"])
                textbox(slide, px0 + slot * i, top - CB_VAL_H, slot, 14, txt,
                        pt=PT["micro"], family=F_DISPLAY, bold=True,
                        color=s_.get("color", "ink2"), align="c", line_pct=1.0,
                        wrap=False)

    # cizgiler
    for s_ in series:
        if s_.get("kind") != "line":
            continue
        side = s_.get("axis", "left")
        if side not in ax:
            continue
        col = C.get(s_.get("color", "coral"), s_.get("color", "coral"))
        # None degeri bosluktur: cizgi orada kesilir, nokta basilmaz. Bir seriyi
        # renk kirilimiyla ikiye bolerken (gecis oncesi / sonrasi) kullanilir.
        pts = []
        for i, v in enumerate(s_["data"][:n]):
            vx = px0 + slot * i + slot / 2
            pts.append(None if v is None else (i, vx, ypos(side, float(v))))
        parca, cari = [], []
        for p in pts:
            if p is None:
                if len(cari) > 1:
                    parca.append(cari)
                cari = []
            else:
                cari.append(p)
        if len(cari) > 1:
            parca.append(cari)
        for seg in parca:
            ff = slide.shapes.build_freeform(px(seg[0][1]), px(seg[0][2]))
            ff.add_line_segments([(px(a), px(bb)) for _, a, bb in seg[1:]], close=False)
            shp = ff.convert_to_shape()
            shp.fill.background()
            shp.line.color.rgb = RGBColor.from_string(col)
            shp.line.width = Pt(2.25)
            shp.shadow.inherit = False
        lbls = s_.get("labels_text")
        for i, vx, vy in [p for p in pts if p]:
            d = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(vx - 4), px(vy - 4),
                                       px(8), px(8))
            d.fill.solid()
            d.fill.fore_color.rgb = RGBColor.from_string(col)
            _no_line(d)
            d.shadow.inherit = False
            if s_.get("labels") == "above":
                v = float(s_["data"][i] or 0)
                txt = lbls[i] if lbls and i < len(lbls) else _fmt_val(v, ax[side]["fmt"])
                textbox(slide, vx - slot / 2, vy - 19, slot, 14, txt,
                        pt=PT["micro"], family=F_DISPLAY, bold=True,
                        color=s_.get("color", "coral"), align="c", line_pct=1.0,
                        wrap=False)

    # kategori etiketleri
    for i, cat in enumerate(cats):
        textbox(slide, px0 + slot * i, plot_bot + 5, slot, 16, str(cat),
                pt=PT["micro"], color="ink2", align="c", line_pct=1.0, wrap=False)

    _label_fit_check(ctx, idx, "combo", slot, [str(c) for c in cats], [])

    ctx.boxes.append((idx, "combo", x, y0, w, h))
    return h



# ----------------------------------------------------------------------------
# Blok: panel kartlari / not kutusu / duz metin / gorsel
# ----------------------------------------------------------------------------

def block_panels(slide, b, x, y, w, ctx, idx):
    """b = {type:"panels", cols?:3, items:[{title, sub?, lines:[..]}], h?:auto}"""
    items = b.get("items") or []
    cols = b.get("cols") or min(3, max(1, len(items)))
    gap = b.get("gap", 18)
    cw = (w - gap * (cols - 1)) / cols
    pad = 18
    pt_ = b.get("font_pt", PT["sm"])

    heights = []
    for it in items:
        hh = pad
        hh += PT["h4"] * PX_PER_PT * 1.25 + 6
        if it.get("sub"):
            hh += len(wrap_lines(plain(it["sub"]), cw - pad * 2, PT["micro"])) \
                  * PT["micro"] * PX_PER_PT * 1.4 + 8
        for ln in it.get("lines") or []:
            hh += len(wrap_lines(plain(ln), cw - pad * 2 - 14, pt_)) \
                  * pt_ * PX_PER_PT * 1.45 + 6
        heights.append(hh + pad)
    ch = b.get("h") or (max(heights) if heights else 100)

    for i, it in enumerate(items):
        cx = x + (i % cols) * (cw + gap)
        cy = y + (i // cols) * (ch + gap)
        rect(slide, cx, cy, cw, ch, fill=b.get("fill", "white"), radius=16,
             line="line", line_w=0.75)
        ty = cy + pad
        textbox(slide, cx + pad, ty, cw - pad * 2, PT["h4"] * PX_PER_PT * 1.25,
                plain(it.get("title", "")), pt=PT["h4"], family=F_DISPLAY,
                bold=True, color=it.get("color", "ink"), line_pct=1.25)
        ty += PT["h4"] * PX_PER_PT * 1.25 + 6
        if it.get("sub"):
            hh = len(wrap_lines(plain(it["sub"]), cw - pad * 2, PT["micro"])) \
                 * PT["micro"] * PX_PER_PT * 1.4
            textbox(slide, cx + pad, ty, cw - pad * 2, hh, it["sub"],
                    pt=PT["micro"], color="ink3", line_pct=1.4)
            ty += hh + 8
        for ln in it.get("lines") or []:
            hh = len(wrap_lines(plain(ln), cw - pad * 2 - 14, pt_)) \
                 * pt_ * PX_PER_PT * 1.45
            textbox(slide, cx + pad, ty, 12, hh, "•", pt=pt_, color="coral",
                    line_pct=1.45)
            textbox(slide, cx + pad + 14, ty, cw - pad * 2 - 14, hh, ln, pt=pt_,
                    color="ink", line_pct=1.45)
            ty += hh + 6

    rows = (len(items) + cols - 1) // cols
    h = rows * ch + (rows - 1) * gap
    ctx.boxes.append((idx, "panels", x, y, w, h))
    return h


def block_note(slide, b, x, y, w, ctx, idx):
    """b = {type:"note", label:"YONTEM", text:"..", fill?:"mint"}"""
    label = plain(b.get("label", "NOT")).upper()
    pad = 16
    pt_ = min(b.get("font_pt", PT["body"]), BODY_PT_MAX)
    lines = wrap_lines(plain(b.get("text", "")), w - pad * 2, pt_)
    th = len(lines) * pt_ * PX_PER_PT * 1.5
    h = pad + PT["xs"] * PX_PER_PT * 1.2 + 8 + th + pad
    rect(slide, x, y, w, h, fill=b.get("fill", "mint"), radius=12)
    textbox(slide, x + pad, y + pad, w - pad * 2, PT["xs"] * PX_PER_PT * 1.2,
            label, pt=PT["xs"], family=F_DISPLAY, bold=True, color="coral_deep",
            line_pct=1.2)
    textbox(slide, x + pad, y + pad + PT["xs"] * PX_PER_PT * 1.2 + 8,
            w - pad * 2, th, b.get("text", ""), pt=pt_, color="ink", line_pct=1.5)
    ctx.boxes.append((idx, "note", x, y, w, h))
    return h


def block_text(slide, b, x, y, w, ctx, idx):
    pt_ = min(b.get("font_pt", PT["body"]), BODY_PT_MAX)
    paras = b.get("paras") or ([b["text"]] if b.get("text") else [])
    total = 0
    rendered = []
    for p in paras:
        ls = wrap_lines(plain(p), w, pt_)
        total += len(ls) * pt_ * PX_PER_PT * 1.55 + 8
        rendered.append(parse_runs(p, b.get("color", "ink")))
    if b.get("title"):
        textbox(slide, x, y, w, PT["h4"] * PX_PER_PT * 1.3, b["title"],
                pt=PT["h4"], family=F_DISPLAY, bold=True, line_pct=1.2)
        y += PT["h4"] * PX_PER_PT * 1.3 + 8
        total += PT["h4"] * PX_PER_PT * 1.3 + 8
    textbox(slide, x, y, w, total, rendered, pt=pt_, color=b.get("color", "ink"),
            line_pct=1.55, space_after=6)
    ctx.boxes.append((idx, "text", x, y, w, total))
    return total


def block_image(slide, b, x, y, w, ctx, idx):
    p = ctx.resolve(b.get("path"))
    if not p or not os.path.exists(p):
        ctx.warn(f"S{idx}: gorsel bulunamadi -> {b.get('path')}")
        return 0
    ph = b.get("h")
    pic = picture(slide, p, x, y, w=b.get("w", w), h=ph)
    h = (pic.height / PX) if pic else 0
    ctx.boxes.append((idx, "image", x, y, w, h))
    return h


BLOCKS = {
    "table": block_table, "insights": block_insights, "kpi": block_kpi,
    "bar": block_bar, "line": block_line, "combo": block_combo,
    "panels": block_panels,
    "note": block_note, "text": block_text, "image": block_image,
}

# Her blok tipinin icerigi hangi alandan okudugu. Yanlis alan adi verildiginde
# (ornegin note'a "text" yerine "body") blok sessizce bos ciziliyor, yuksekligi
# dogru oldugu icin tasma denetimine de takilmiyordu. Gercek olay: uc destede
# panels gövdeleri ve note metinleri bos yayinlandi (bkz. tuzaklar 3.6d).
# Ic listeler alternatif alanlari gosterir: en az biri dolu olmali.
BLOK_ICERIK = {
    "table": [["rows"]],
    "insights": [["items"]],
    "kpi": [["cards"]],
    "bar": [["cats"], ["series"]],
    "line": [["cats"], ["series"]],
    "combo": [["cats"], ["series"]],
    "panels": [["items"]],
    "note": [["text"]],
    "text": [["paras", "text"]],
    "image": [["src"]],
}
# panels/kpi gibi liste tasiyan bloklarda her ogenin icerik alani
OGE_ICERIK = {"panels": "lines", "kpi": "value"}
# icerik alani sanilip yanlislikla yazilan yaygin adlar
TAKMA_AD = {"body", "content", "desc", "description", "metin", "govde", "icerik"}


def blok_icerik_denetimi(b, ctx, idx):
    """Blogun icerik alanlari dolu mu? Bos ya da yanlis adlandirilmissa uyarir."""
    tip = b.get("type")
    for secenek in BLOK_ICERIK.get(tip, []):
        if not any(b.get(a) for a in secenek):
            ad = " ya da ".join(f"'{a}'" for a in secenek)
            ctx.warn(f"BOS BLOK S{idx}: '{tip}' blogunda {ad} alani bos ya da yok - "
                     f"blok gorunur icerik uretmeyecek")
    oge_alan = OGE_ICERIK.get(tip)
    if oge_alan:
        for i, oge in enumerate(b.get("items") or b.get("cards") or []):
            if isinstance(oge, dict) and not oge.get(oge_alan):
                ctx.warn(f"BOS BLOK S{idx}: '{tip}' blogunun {i+1}. ogesinde "
                         f"'{oge_alan}' alani bos ya da yok")
    # yanlis alan adi: icerik gibi duran ama taninmayan anahtar
    bilinmeyen = TAKMA_AD & set(b)
    if bilinmeyen:
        beklenen = " / ".join(" ya da ".join(x) for x in BLOK_ICERIK.get(tip, [])) or "?"
        ctx.warn(f"BILINMEYEN ALAN S{idx}: '{tip}' blogunda {sorted(bilinmeyen)} "
                 f"anahtari kullanilmis - bu tip '{beklenen}' bekliyor, deger yok sayilacak")
    for oge in (b.get("items") or []):
        if isinstance(oge, dict) and (TAKMA_AD & set(oge)):
            ctx.warn(f"BILINMEYEN ALAN S{idx}: '{tip}' oge icinde "
                     f"{sorted(TAKMA_AD & set(oge))} kullanilmis - "
                     f"beklenen alan '{oge_alan or '?'}'")


# ----------------------------------------------------------------------------
# Slayt tipleri
# ----------------------------------------------------------------------------

def s_cover(slide, spec, ctx, idx):
    """
    VitrA destesindeki kapak: coral zemin, sol kenara yaslanmis soluk big-O,
    ortalanmis marka + deste tipi basligi, hemen altinda neredeyse ayni puntoda
    donem satiri, alt-ortada wordmark. Metin rengi saf beyaz degil paper (#FEFEF7).
    """
    set_bg(slide, spec.get("bg", "coral"))
    art = ctx.slide_img("cover-art-front.png")
    if os.path.exists(art):
        picture(slide, art, 0, 0, w=COVER_ART_W, h=STAGE_H)

    avail = STAGE_W - 120
    title = plain(spec.get("title", ""))
    lines = spec.get("title_lines") or ([title] if title else [])
    tpt = spec.get("title_pt", COVER_TITLE_PT)
    for ln in lines:
        tpt = min(tpt, fit_pt(ln, avail, tpt, 24, F_DISPLAY, weight=COVER_TITLE_W))
    lh = tpt * PX_PER_PT * 1.07

    y = spec.get("title_y", COVER_TITLE_Y)
    for ln in lines:
        textbox(slide, 60, y, avail, lh, ln, pt=tpt, family=F_DISPLAY, bold=True,
                color="paper", align="c", line_pct=1.07, wrap=False)
        y += lh

    if spec.get("subtitle"):
        spt = spec.get("subtitle_pt", COVER_SUB_PT)
        sub = plain(spec["subtitle"])
        spt = fit_pt(sub, avail, spt, 16, F_DISPLAY, weight=COVER_SUB_W)
        y += COVER_SUB_GAP
        textbox(slide, 60, y, avail, spt * PX_PER_PT * 1.05, sub, pt=spt,
                family=F_DISPLAY, bold=True, color="paper", align="c",
                line_pct=1.05, wrap=False)

    wm = ctx.slide_img("wordmark-cover.png")
    if os.path.exists(wm):
        x, wy, ww, wh = COVER_WM
        picture(slide, wm, x, wy, w=ww, h=wh)


def s_agenda(slide, spec, ctx, idx):
    """
    VitrA ajandasi: kirik-beyaz zemin, sol yarida coral panel (sag kenarlari
    yuvarlatilmis) + uzerinde soluk big-O, panelde ortalanmis tek satir baslik,
    sol altta beyaz logo. Sagda numarali liste: numara ustte normal, etiket
    altta kalin, teal renkte; blok dikeyde ortalanir.
    """
    set_bg(slide, spec.get("bg", "paper_bg"))
    panel = ctx.slide_img("agenda-panel.png")
    if os.path.exists(panel):
        picture(slide, panel, -8, 0, w=AGENDA_PANEL_W, h=STAGE_H)
    else:
        rect(slide, 0, 0, AGENDA_PANEL_W, STAGE_H, fill="coral")
    art = ctx.slide_img("cover-art-front.png")
    if os.path.exists(art):
        picture(slide, art, -8, 0, w=COVER_ART_W, h=STAGE_H)

    if spec.get("kicker"):
        textbox(slide, AGENDA_EYEBROW_XY[0], AGENDA_EYEBROW_XY[1],
                AGENDA_PANEL_W - 60, 18, plain(spec["kicker"]).upper(),
                pt=10, family=F_BODY, color="paper", line_pct=1.0)

    lines = spec.get("title_lines") or [plain(spec.get("title", "SUNUM AKIŞI"))]
    tpt = spec.get("title_pt", AGENDA_TITLE_PT)
    for ln in lines:
        tpt = min(tpt, fit_pt(ln, AGENDA_PANEL_W - 60, tpt, 20, F_DISPLAY,
                              bold=False, weight=AGENDA_TITLE_W))
    lh = tpt * PX_PER_PT * 1.15
    ty = spec.get("title_y", AGENDA_TITLE_Y) - (len(lines) - 1) * lh / 2
    for ln in lines:
        textbox(slide, -3, ty, AGENDA_PANEL_W + 6, lh, ln, pt=tpt,
                family=F_DISPLAY, bold=False, color="paper", align="c",
                line_pct=1.15, wrap=False)
        ty += lh

    lg = ctx.slide_img("agenda-logo-inner.png")
    if os.path.exists(lg):
        x, ly, lw_, lh_ = AGENDA_LOGO
        picture(slide, lg, x, ly, w=lw_, h=lh_)

    # sag kolon: numarali liste, dikeyde ortalanmis
    items = spec.get("items") or []
    ipt = spec.get("item_pt", AGENDA_ITEM_PT)
    nlh = ipt * PX_PER_PT * AGENDA_ITEM_LH
    rows = []
    for i, it in enumerate(items):
        no = (it.get("no") if isinstance(it, dict) else None) or f"{i+1:02d}"
        label = plain(it.get("label") if isinstance(it, dict) else str(it))
        wrapped = wrap_lines(label, AGENDA_LIST_W, ipt, F_DISPLAY, True)
        rows.append((str(no), wrapped))
    total = (sum(nlh * (1 + len(w)) + AGENDA_ITEM_NUM_GAP for _n, w in rows)
             + AGENDA_ITEM_GAP * max(0, len(rows) - 1))
    y = (STAGE_H - total) / 2
    if y < 40:
        y = 40
        ctx.warn(f"S{idx}: ajanda listesi {len(rows)} madde ile dikeyde sigmiyor, "
                 f"ust sinira yaslandi - madde sayisini azaltmak veya item_pt "
                 f"kucultmek degerlendirilebilir")
    for no, wrapped in rows:
        textbox(slide, AGENDA_LIST_X, y, AGENDA_LIST_W, nlh, no, pt=ipt,
                family=F_DISPLAY, bold=False, color="ink3", line_pct=AGENDA_ITEM_LH,
                wrap=False)
        y += nlh + AGENDA_ITEM_NUM_GAP
        textbox(slide, AGENDA_LIST_X, y, AGENDA_LIST_W, nlh * len(wrapped),
                [parse_runs(ln, "ink") for ln in wrapped], pt=ipt,
                family=F_DISPLAY, bold=True, color="ink",
                line_pct=AGENDA_ITEM_LH, wrap=False)
        y += nlh * len(wrapped) + AGENDA_ITEM_GAP


def separator_layout(spec):
    """
    Ayrac yerlesimi - VitrA destesinden olculmus degerlerle.

    Numeral SABIT konumdadir: 200pt ExtraBold, sep_num renginde, yatay merkezi
    x=146. VitrA'da numeral bir filigran gibi davraniyor; basligin genisligine
    gore yer degistirmiyor ve gerektiginde basligin arkasinda kaliyor. Punto da
    sabittir - numerali basliga gore olceklendirmek ayni destede farkli boyutta
    numeraller uretiyordu.

    Baslik sayfa ortasinda, 37pt ExtraBold, paper renginde. SEP_TITLE_MAX_W'yi
    asarsa alt satira kayar; accent cizgiler (coral) blok yuksekligine gore
    simetrik olarak acilir ve blok dikeyde ortalanir.

    HTML onizleme de bu fonksiyonu kullanir; satir kirilmasi tek yerde belirlenir.
    """
    title = plain(spec.get("title", ""))
    tpt = spec.get("title_pt", SEP_TITLE_PT)
    num = str(spec.get("no", "") or "")
    npt = spec.get("no_pt", SEP_NUM_PT)
    nw = text_w(num, npt, F_DISPLAY, True, SEP_NUM_W) if num else 0.0

    max_tw = spec.get("title_max_w", SEP_TITLE_MAX_W)
    lines = wrap_lines(title, max_tw, tpt, F_DISPLAY, True, weight=SEP_TITLE_W)
    tw = max((text_w(ln, tpt, F_DISPLAY, True, SEP_TITLE_W) for ln in lines),
             default=0.0)

    lh = tpt * PX_PER_PT * 1.08
    th = len(lines) * lh
    cy = STAGE_H / 2
    total_h = SEP_ACC_H + SEP_ACC_GAP + th + SEP_ACC_GAP + SEP_ACC_H
    top = cy - total_h / 2

    return dict(
        title=title, lines=lines, tpt=tpt, tw=tw, th=th, lh=lh,
        tx=(STAGE_W - tw) / 2, num=num, npt=npt, nw=nw,
        max_tw=max_tw, cy=cy,
        num_x=SEP_NUM_CX - nw / 2, num_h=npt * PX_PER_PT,
        acc_top_y=top,
        title_y=top + SEP_ACC_H + SEP_ACC_GAP,
        acc_bot_y=top + SEP_ACC_H + SEP_ACC_GAP + th + SEP_ACC_GAP,
    )


def s_separator(slide, spec, ctx, idx):
    set_bg(slide, spec.get("bg", "teal"))
    L = separator_layout(spec)

    if len(L["lines"]) > 2:
        ctx.warn(f"S{idx}: ayrac basligi {len(L['lines'])} satira sarildi - iki "
                 f"satir tasarimin rahat siniri. Basligi kisaltmak onerilir "
                 f"('{L['title']}')")

    # numeral once cizilir: filigran olarak basligin arkasinda kalir
    if L["num"]:
        textbox(slide, L["num_x"] - 20, L["cy"] - L["num_h"] * 0.52,
                L["nw"] + 40, L["num_h"], L["num"], pt=L["npt"],
                family=F_DISPLAY, bold=True, color="sep_num", align="c",
                line_pct=0.9, wrap=False)

    accx = (STAGE_W - SEP_ACC_W) / 2
    rect(slide, accx, L["acc_top_y"], SEP_ACC_W, SEP_ACC_H, fill="coral", radius=3)
    textbox(slide, L["tx"] - 30, L["title_y"], L["tw"] + 60, L["th"],
            [parse_runs(ln, "paper") for ln in L["lines"]], pt=L["tpt"],
            family=F_DISPLAY, bold=True, color="paper", align="c",
            line_pct=1.08, wrap=False)
    rect(slide, accx, L["acc_bot_y"], SEP_ACC_W, SEP_ACC_H, fill="coral", radius=3)


def s_closing(slide, spec, ctx, idx):
    set_bg(slide, spec.get("bg", "teal"))
    big = ctx.logo("inbound-big-o-white.png")
    if os.path.exists(big):
        pic = picture(slide, big, -300, 200, h=1000)
        if pic:
            _fade(pic, 10000)
    title = plain(spec.get("title", "Teşekkürler"))
    tpt = spec.get("title_pt", 56)
    while tpt > 28 and text_w(title, tpt, F_DISPLAY, True) > STAGE_W - 200:
        tpt -= 1
    textbox(slide, 80, STAGE_H / 2 - 60, STAGE_W - 160, tpt * PX_PER_PT * 1.1,
            title, pt=tpt, family=F_DISPLAY, bold=True, color="white",
            align="c", line_pct=1.1, wrap=False)
    if spec.get("subtitle"):
        textbox(slide, 160, STAGE_H / 2 + 10, STAGE_W - 320, 60, spec["subtitle"],
                pt=PT["h4"], color="white", align="c", line_pct=1.4)
    wm = ctx.logo("inbound-wordmark-white.png")
    if os.path.exists(wm):
        pic = picture(slide, wm, 0, STAGE_H - 94, h=30)
        if pic:
            pic.left = px((STAGE_W - pic.width / PX) / 2)


def s_content(slide, spec, ctx, idx):
    dark = spec.get("bg") in ("teal", "dark")
    if spec.get("bg") and spec["bg"] not in ("white", None):
        set_bg(slide, "teal" if spec["bg"] == "dark" else spec["bg"])
    top = chrome(slide, spec, ctx, idx, dark=dark)
    bottom = footnotes(slide, spec.get("footnotes"), ctx, idx)

    grid = spec.get("grid") or [100]
    gap = spec.get("gap", 36)
    avail = STAGE_W - M_L - M_R
    tot = float(sum(grid))
    xs, ws, cx = [], [], M_L
    for g in grid:
        cw = (avail - gap * (len(grid) - 1)) * g / tot
        xs.append(cx)
        ws.append(cw)
        cx += cw + gap

    cursor = [top] * len(grid)
    auto = 0
    for b in spec.get("blocks") or []:
        fn = BLOCKS.get(b.get("type"))
        if not fn:
            ctx.warn(f"S{idx}: bilinmeyen blok tipi '{b.get('type')}'")
            continue
        blok_icerik_denetimi(b, ctx, idx)
        if b.get("at"):
            x, y, w, _h = b["at"]
            fn(slide, b, x, y, w, ctx, idx)
            continue
        col = b.get("col")
        if col == "full":
            # Tum kolonlarin altinda, tam genislikte. Slayt katalogunda sik gecen
            # "iki tablo yan yana, altta yorum paragrafi" deseni (C30, C33) icin.
            y = max(cursor) + (b.get("mt") or 0)
            used = fn(slide, b, M_L, y, avail, ctx, idx)
            nxt = y + used + b.get("mb", 20)
            cursor = [nxt] * len(cursor)
            if y + used > bottom + 2:
                ctx.warn(f"TASMA S{idx}: '{b.get('type')}' blogu (tam genislik) "
                         f"y={y+used:.0f}px, alt sinir {bottom:.0f}px - "
                         f"{y+used-bottom:.0f}px asiyor")
            continue
        if col is None:
            col = auto % len(grid)
            auto += 1
        col = max(0, min(col, len(grid) - 1))
        y = cursor[col] + (b.get("mt") or 0)
        used = fn(slide, b, xs[col], y, ws[col], ctx, idx)
        cursor[col] = y + used + b.get("mb", 20)
        if y + used > bottom + 2:
            ctx.warn(f"TASMA S{idx}: '{b.get('type')}' blogu (kolon {col}) "
                     f"y={y+used:.0f}px, alt sinir {bottom:.0f}px - "
                     f"{y+used-bottom:.0f}px asiyor")
            ctx.overflowed.add((idx, b.get("type")))


SLIDES = {"cover": s_cover, "agenda": s_agenda, "separator": s_separator,
          "content": s_content, "closing": s_closing}


def _fade(pic, alpha_off):
    """Gorsele saydamlik uygular (big-O dekoratif marka)."""
    try:
        blip = pic._element.blipFill.find(qn("a:blip"))
        from lxml import etree
        el = etree.SubElement(blip, qn("a:alphaModFix"))
        el.set("amt", str(100000 - alpha_off))
    except Exception:
        pass


# ----------------------------------------------------------------------------
# Deste kurulumu
# ----------------------------------------------------------------------------

def _blank_layout(prs):
    """En bos layout'u sec ve mirasli placeholder'lari temizle."""
    best, best_n = prs.slide_layouts[0], 99
    for lay in prs.slide_layouts:
        n = len(lay.placeholders)
        if n < best_n:
            best, best_n = lay, n
    return best


def build(spec, out_path, assets_dir, base_dir, check_only=False):
    prs = Presentation()
    prs.slide_width = px(STAGE_W)
    prs.slide_height = px(STAGE_H)
    layout = _blank_layout(prs)
    ctx = Ctx(assets_dir, base_dir)

    slides_spec = spec.get("slides") or []
    for i, s in enumerate(slides_spec, 1):
        kind = s.get("type", "content")
        fn = SLIDES.get(kind)
        if not fn:
            ctx.warn(f"S{i}: bilinmeyen slayt tipi '{kind}' - atlandi")
            continue
        slide = prs.slides.add_slide(layout)
        for shp in list(slide.shapes):
            if shp.is_placeholder:
                shp._element.getparent().remove(shp._element)
        if kind == "content" and not s.get("bg"):
            set_bg(slide, "white")
        fn(slide, s, ctx, i)
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = plain(s["notes"])

    # sag kenar taramasi (alt sinir taramasi s_content icinde, kolon bilgisiyle)
    for (si, label, x, y, w, h) in ctx.boxes:
        if x + w > STAGE_W - M_R + 2:
            ctx.warn(f"TASMA S{si}: {label} blogu sag kenari asiyor "
                     f"(x+w={x+w:.0f}px, sinir {STAGE_W-M_R}px)")

    if not check_only:
        prs.save(out_path)
    return ctx, len(slides_spec)


def main():
    ap = argparse.ArgumentParser(description="deck.json -> Inbound PPTX")
    ap.add_argument("spec")
    ap.add_argument("-o", "--out", default=None)
    ap.add_argument("--assets", default=None,
                    help="varsayilan: <script>/../assets")
    ap.add_argument("--check", action="store_true", help="dosya yazmadan dogrula")
    a = ap.parse_args()

    here = os.path.dirname(os.path.abspath(__file__))
    assets = a.assets or os.path.normpath(os.path.join(here, "..", "assets"))
    with open(a.spec, encoding="utf-8") as f:
        spec = json.load(f)
    base = os.path.dirname(os.path.abspath(a.spec))
    out = a.out or spec.get("output") or os.path.splitext(a.spec)[0] + ".pptx"
    if not os.path.isabs(out):
        out = os.path.join(base, out)

    ctx, n = build(spec, out, assets, base, check_only=a.check)

    if FONT_LOAD_ERROR:
        print(f"\n!!! FONT UYARISI: gercek font metrigi kullanilamiyor "
              f"({FONT_LOAD_ERROR}).\n    Olcum karakter genisligi yaklasimina "
              f"dustu; tasma tespiti ve baslik auto-shrink guvenilir degil.\n"
              f"    assets/design-system/fonts/ altindaki variable TTF'leri "
              f"kontrol et.\n")
    print(f"{n} slayt {'dogrulandi' if a.check else 'uretildi'}"
          f"{'' if a.check else ' -> ' + out}")
    if ctx.warnings:
        print(f"\n{len(ctx.warnings)} uyari:")
        for w in ctx.warnings:
            print("  -", w)
        return 1
    print("Uyari yok.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
