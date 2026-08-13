#!/usr/bin/env python3
"""
qa_deck.py - Teslim oncesi tek kapi denetim.

Uc katman tarar:
  1. YERLESIM   - inbound_deck.py'nin olcum motorunu check modunda calistirir
                  (tasma, baslik sarmasi, sag kenar asimi).
  2. DIL         - icerik-dili-rehberi kalemleri: em dash, mojibake, emoji,
                  emir kipi, kesin vaat, yasakli keskin kelime, "ss."/"pik",
                  cift bosluk, otomasyon araci adi sizintisi, ic kisit ifadesi.
  3. RAKAM/YAPI  - yuzde formati tutarliligi, ondalik ayirici karisikligi,
                  imkansiz yuzde, kaynak notu eksikligi, ajanda-bolum uyumu,
                  placeholder kalintisi, marka yazim tutarliligi.

Bulgular iki siniftir:
  HATA   - teslim engelleyici, duzeltilmeden cikilmaz.
  UYARI  - insan karari gerektirir (ornegin bilincli istisna olabilir).

Kullanim:
    python3 qa_deck.py deck.json
    python3 qa_deck.py deck.json --pptx cikti.pptx   # uretilmis dosyayi da tara
    python3 qa_deck.py deck.json --json              # makine okunur cikti
"""

from __future__ import annotations

import argparse
import json
import os
import re
import sys
import unicodedata

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from inbound_deck import build, plain  # noqa: E402

# ----------------------------------------------------------------------------
# Dil kurallari (icerik-dili-rehberi Bolum 17 self-check karsiliklari)
# ----------------------------------------------------------------------------

EMOJI = re.compile(
    "[" "\U0001F300-\U0001FAFF" "\U0001F000-\U0001F2FF"
    "☀-➿" "️" "\U0001F1E6-\U0001F1FF" "]",
    flags=re.UNICODE)
# islevsel semboller serbest: ➔ → ↑ ↓ ✓ ▲ · ★
ALLOWED_SYM = set("➔→↑↓✓▲·★—…’‘“”")

MOJIBAKE = [("Ã", "UTF-8 cift kodlama"), ("Ä±", "UTF-8 cift kodlama"),
            ("Â·", "UTF-8 cift kodlama"), ("﻿", "ZWNBSP"),
            ("​", "zero-width space")]

# emir kipi: 2. cokul emir ekleri. Yanlis pozitifi azaltmak icin fiil govdesi
# uzunlugu ve yaygin isim istisnalari filtreleniyor.
IMPERATIVE = re.compile(
    r"\b\w{3,}(?:ınız|iniz|unuz|ünüz)\b|"
    r"\b(?:ekle|kaldır|düzelt|güncelle|yap|kur|aç|kapat|artır|azalt|optimize\s+et)"
    r"(?:yin|yiniz|in|iniz|ın|ınız|un|unuz|ün|ünüz)\b",
    re.IGNORECASE)
IMP_SAFE = {"inceleyebilirsiniz", "görebilirsiniz", "izleyebilirsiniz",
            "bulabilirsiniz", "ulaşabilirsiniz", "erişebilirsiniz"}

PROMISE = [r"\bgarantil", r"\bgaranti\b", r"\bmutlaka\b", r"\bkesinlikle\b",
           r"\bkesin\b", r"\bhemen\b", r"\bacilen\b", r"\bçok kolay\b",
           r"\bşüphesiz\b"]

HARSH = {
    "kötü": "geliştirme alanı barındıran", "berbat": "geliştirme alanı",
    "başarısız": "hedeflenen sonuca ulaşmayan", "yanlış": "gözden geçirilebilir",
    "zayıf": "büyüme potansiyeli taşıyan", "ciddi": "dikkate değer",
    "patlama": "yüksek artış", "dramatik": "belirgin",
    "büyüme motoru": "büyüme kanalı", "kat kat": "oran veya Nx çarpan",
    "en zayıf halka": "en az işlenen başlık", "ters profil": "tersine dönen örüntü",
    "açık ara": "belirgin biçimde", "cannibalization": "fırsat kaybı",
    "para tuzağı": "ek süre maliyeti", "domine": "belirgin biçimde önde",
    "ele geçir": "nötr karşılık", "mercek": "perspektif",
    "doğrulanmış fırsat": "fırsat", "panorama": "nötr karşılık",
    # Erdogan geri bildirimi: "somut/kanit" tipi pekistirme ifadeleri
    # kullanilmiyor - bulgu kendi basina durur.
    "somut": "ifadeyi kaldir, bulguyu dogrudan ver",
    "kanıt": "ornek / ekran goruntusu / dayanak",
}
# teknik adlandirma olarak serbest kullanimlar
HARSH_CONTEXT_OK = re.compile(
    r"(canonical|404|301|302|500|schema|indexleme|crawl|redirect|status\s*code)"
    r"\s*(hata|sorun)", re.IGNORECASE)
JUDGMENT = ["hata", "sorun"]

# Rehberin notr teknik fiil sozlugu ("dusus, gerileme, daralma, ... zayiflama")
# yasak sifatin turevlerini mesru kiliyor: "zayiflamistir" kabul, "zayif" degil.
HARSH_ALLOWED_FORMS = {
    "zayıf": re.compile(r"zayıfla"),
    "hata":  re.compile(r"hatas[ıi]z|hatal[ıi] kosu"),
}

TOOL_LEAK = ["claude", "claude code", "mcp", "playwright", "skill",
             "anthropic", "python-pptx", "pptxgenjs", "matplotlib",
             "dataforseo", "screaming frog", "markitdown"]
INTERNAL = ["veri çekilemedi", "ölçülemiyor", "ölçülemedi", "erişemedim",
            "erişilemedi", "çekemedim", "api limiti", "rate limit",
            "sonraki aşamada eklenecek", "eklenecek", "in progress",
            "not started", "blocked", "(shared)", "todo", "tbd", "lorem",
            "section title", "xxx", "placeholder"]

BANNED_TERMS = {"ss.": "session", "pik": "peak", "atıf": "mention",
                "makine okunur": "machine-readable", "üst fold": "above the fold",
                "3th party": "3rd party", "GA360": "GA4",
                "çapa metni": "anchor text", "market kümesi": "Market Cluster"}

PCT_OK = re.compile(r"[+\-]?%\d")               # dogru: %18, +%6.9, -%37
PCT_BAD = re.compile(r"\d[\d.,]*\s*%")          # yanlis: 34.8%, 2,75%
PCT_SIGN_BAD = re.compile(r"%\s*[+\-]\s*\d")    # yanlis: %-3.3, %+12
DEC_COMMA = re.compile(r"\d,\d")                # yanlis ondalik: 2,75
THOUSAND_DOT = re.compile(r"\b\d{1,3}(?:\.\d{3})+\b")   # dogru binlik: 17.637
IMPOSSIBLE = re.compile(r"-\s*%\s*(1[0-9]{2,}|[2-9][0-9]{2,})(?:\.\d+)?\s*"
                        r"(?:düşüş|azalma|gerileme|daralma|kayıp)", re.IGNORECASE)
BIG_PCT = re.compile(r"[+\-]?%\s*(\d{4,})")


class Report:
    def __init__(self):
        self.items = []

    def add(self, level, where, rule, detail, fix=""):
        self.items.append(dict(level=level, where=where, rule=rule,
                               detail=detail, fix=fix))

    def err(self, *a, **k):
        self.add("HATA", *a, **k)

    def warn(self, *a, **k):
        self.add("UYARI", *a, **k)

    @property
    def errors(self):
        return [i for i in self.items if i["level"] == "HATA"]

    @property
    def warnings(self):
        return [i for i in self.items if i["level"] == "UYARI"]


# ----------------------------------------------------------------------------
# Metin toplama
# ----------------------------------------------------------------------------

def collect(spec):
    """(konum, metin) ciftleri - kaynak alanini isaretleyerek."""
    out = []

    def walk(node, where):
        if isinstance(node, str):
            out.append((where, node))
        elif isinstance(node, list):
            for i, v in enumerate(node):
                walk(v, where)
        elif isinstance(node, dict):
            for k, v in node.items():
                if k.startswith("_") or k in {"path", "output", "color", "fill",
                                              "accent", "align", "type"}:
                    continue
                walk(v, f"{where}.{k}" if where else k)

    for i, s in enumerate(spec.get("slides") or [], 1):
        walk(s, f"S{i:02d}")
    return out


# ----------------------------------------------------------------------------
# Katman 2: dil
# ----------------------------------------------------------------------------

def alinti_mi(t):
    """Alan bastan sona bir alintidan mi ibaret?

    Sosyal dinleme ve AI gorunurluk slaytlarinda kullanici/model yanitlari
    BIREBIR aktarilir (icerik-dili-rehberi 14.6). Bu metinlerde emir kipi,
    marka yaziminin farkli hali ve keskin kelime bulunabilir; bunlar destenin
    kendi sesi degildir ve duzeltilmez. Yalnizca tirnakla acilip kapanan tam
    alintilar muaf tutulur - alintinin icine yerlestirilmis yorum cumlesi
    denetimden kacamaz.
    """
    t = t.strip()
    ac = "\"“‘'"
    kapa = "\"”’'"
    return len(t) > 20 and t[0] in ac and t[-1] in kapa


def check_language(spec, rep):
    for where, raw in collect(spec):
        t = plain(raw)
        low = t.lower()
        if alinti_mi(t):
            continue

        if "—" in t or "–" in t:
            rep.err(where, "em dash", f"'{t[:70]}'",
                    "kısa tire '-', iki nokta veya virgülle değiştir")

        for bad, why in MOJIBAKE:
            if bad in t:
                rep.err(where, "mojibake", f"{why}: {repr(bad)} bulundu",
                        "UTF-8 temizliği yap (Python ile, sed/perl ile değil)")

        for ch in EMOJI.findall(t):
            if ch not in ALLOWED_SYM:
                rep.err(where, "emoji", f"'{ch}' ({unicodedata.name(ch,'?')})",
                        "kaldır; ✓ ▲ ↑ ↓ ➔ işlevsel sembolleri serbest")

        for m in IMPERATIVE.finditer(t):
            if m.group(0).lower() not in IMP_SAFE:
                rep.err(where, "emir kipi", f"'{m.group(0)}'",
                        "'-ebilir / önerilebilir / değerlendirilebilir'e çevir")

        for pat in PROMISE:
            m = re.search(pat, low)
            if m:
                rep.err(where, "kesin vaat", f"'{m.group(0)}' -> '{t[:60]}'",
                        "'potansiyel taşımaktadır / değerlendirilebilir'e çevir")

        for bad, good in HARSH.items():
            allow = HARSH_ALLOWED_FORMS.get(bad)
            hit = False
            for m in re.finditer(r"\b" + re.escape(bad), low):
                if allow and allow.match(low, m.start()):
                    continue
                hit = True
                break
            if hit:
                rep.err(where, "keskin/abartı kelime", f"'{bad}'",
                        f"'{good}' kullan")

        for j in JUDGMENT:
            for m in re.finditer(r"\b" + j + r"\w*", low):
                seg = t[max(0, m.start() - 30):m.end() + 10]
                if not HARSH_CONTEXT_OK.search(seg):
                    rep.warn(where, "yargı kelimesi", f"'{m.group(0)}' -> '{seg}'",
                             "teknik adlandırma değilse 'bulgu / tespit / "
                             "ele alınabilecek nokta'ya çevir")

        for tool in TOOL_LEAK:
            if re.search(r"\b" + re.escape(tool) + r"\b", low):
                rep.err(where, "otomasyon aracı sızıntısı", f"'{tool}'",
                        "müşteri çıktısında görünmez; jenerikleştir veya kaldır")

        for ic in INTERNAL:
            if ic in low:
                rep.err(where, "iç kısıt / placeholder", f"'{ic}'",
                        "rapordan çıkar, chat üzerinden sunacak kişiye ilet")

        for bad, good in BANNED_TERMS.items():
            if re.search(r"(?<![\w])" + re.escape(bad) + r"(?![\w])", t,
                         re.IGNORECASE):
                rep.err(where, "terim", f"'{bad}'", f"'{good}' kullan")

        if "  " in t:
            rep.warn(where, "çift boşluk", f"'{t[:70]}'", "tek boşluğa indir")
        if re.search(r"[a-zçğıöşü]{2,}\s+ile\s*$", t.strip(), re.IGNORECASE):
            rep.warn(where, "yarım cümle", f"'{t[-60:]}'", "cümleyi tamamla")
        if t.count("!") and not re.search(r"(uyarı|not|dikkat)", low):
            rep.warn(where, "gereksiz ünlem", f"'{t[:60]}'",
                     "gerçek not/uyarı taşımıyorsa kaldır")


# ----------------------------------------------------------------------------
# Katman 3: rakam ve yapi
# ----------------------------------------------------------------------------

def check_numbers(spec, rep):
    for where, raw in collect(spec):
        t = plain(raw)

        for m in PCT_BAD.finditer(t):
            seg = m.group(0)
            if "p" in seg or PCT_OK.search(seg):
                continue
            rep.err(where, "yüzde formatı", f"'{seg}'",
                    "% işareti sayının önünde: '%18', '+%6.9', '-%37'")
        for m in PCT_SIGN_BAD.finditer(t):
            rep.err(where, "yüzde formatı", f"'{m.group(0)}'",
                    "işaret %'den önce gelir: '-%3.3'")
        for m in DEC_COMMA.finditer(t):
            if THOUSAND_DOT.search(t):
                continue
            rep.err(where, "ondalık ayırıcı", f"'{m.group(0)}'",
                    "ondalık ayırıcı her yerde nokta")
        for m in IMPOSSIBLE.finditer(t):
            rep.err(where, "imkansız yüzde", f"'{m.group(0)}'",
                    "düşüş %100'ü aşamaz; mutlak değer veya çarpan kullan")
        for m in BIG_PCT.finditer(t):
            if not re.search(r"\(\s*[+\-]?[\d.,]+\s*(K|M|click|session|adet)?",
                             t[m.end():m.end() + 40], re.IGNORECASE):
                rep.warn(where, "bin üzeri yüzde", f"'%{m.group(1)}'",
                         "yanına mutlak değeri de yaz: '+%1362 (+395 click)'")
        if re.search(r"\+%100\b", t):
            rep.warn(where, "0 tabanlı delta", "'+%100'",
                     "önceki dönem 0 ise 'yeni' etiketi kullan")


def check_structure(spec, rep):
    slides = spec.get("slides") or []
    agenda_items, sep_nos, titles = [], [], []

    for i, s in enumerate(slides, 1):
        t = s.get("type", "content")
        if t == "agenda":
            for it in s.get("items") or []:
                lab = it.get("label") if isinstance(it, dict) else str(it)
                no = (it.get("no") if isinstance(it, dict) else None)
                agenda_items.append((no, plain(lab)))
        if t == "separator":
            sep_nos.append((i, str(s.get("no") or ""), plain(s.get("title", ""))))
        if t == "content":
            titles.append(plain(s.get("title", "")))
            has_data = any(b.get("type") in {"table", "bar", "line", "kpi"}
                           for b in (s.get("blocks") or []))
            if has_data and not s.get("source"):
                rep.err(f"S{i:02d}", "kaynak notu",
                        "veri bloğu var ancak 'source' alanı yok",
                        "her veri slaytına 'Kaynak: <sistem> - <property> · <dönem>' ekle")
            if not s.get("subtitle"):
                rep.warn(f"S{i:02d}", "dönem beyanı", "subtitle yok",
                         "dönem alt başlıkta parantezle beyan edilir")
            for b in s.get("blocks") or []:
                if b.get("type") == "table":
                    rows = b.get("rows") or []
                    ncol = len(b.get("head") or [])
                    for ri, r in enumerate(rows):
                        if ncol and len(r) != ncol:
                            rep.err(f"S{i:02d}", "tablo kolon sayısı",
                                    f"satır {ri}: {len(r)} hücre, başlık {ncol}",
                                    "satır uzunluklarını başlıkla eşitle")
            # yorumsuz tablo kontrolu
            types = [b.get("type") for b in (s.get("blocks") or [])]
            if ("table" in types or "bar" in types or "line" in types) \
                    and not ({"insights", "text", "note"} & set(types)):
                rep.err(f"S{i:02d}", "yorumsuz tablo",
                        "veri bloğu var, insights/text yorum katmanı yok",
                        "her tabloya ➔ yorum katmanı eklenir (müşteri sürümü kuralı)")

    if agenda_items:
        if len(agenda_items) != len(sep_nos):
            rep.err("YAPI", "ajanda-bölüm uyumu",
                    f"ajandada {len(agenda_items)} madde, destede "
                    f"{len(sep_nos)} bölüm ayracı",
                    "sayıları eşitle (bilinen teslim hatası)")
        a_no = [n for n, _ in agenda_items if n]
        s_no = [n for _, n, _ in sep_nos if n]
        if len(set(s_no)) != len(s_no):
            rep.err("YAPI", "tekrarlanan bölüm numarası",
                    f"{s_no}", "her bölüm numarası bir kez kullanılır")
        for (an, al), (si, sn, st) in zip(agenda_items, sep_nos):
            if an and sn and an != sn:
                rep.err("YAPI", "ajanda-ayraç numara farkı",
                        f"ajanda '{an}' vs ayraç '{sn}' ({st})", "numaraları eşle")
            if al and st and al.lower() != st.lower():
                rep.warn("YAPI", "ajanda-ayraç başlık farkı",
                         f"'{al}' vs '{st}'", "aynı adlandırmayı kullan")

    # Marka yazim tutarliligi. Yalnizca destenin KENDI sesinde aranir:
    # tablo satirlarindaki query metinleri ("flormar maskara"), domain adlari
    # ("flormar.com.tr") ve URL yollari markanin kucuk harfli halini mesru
    # olarak icerir; bunlari varyant saymak yanlis pozitif uretiyordu.
    brand = (spec.get("meta") or {}).get("brand")
    if brand:
        variants = set()
        for where, raw in collect(spec):
            if ".rows" in where or ".head" in where:
                continue                      # veri hucreleri: query/URL/domain
            if where.endswith(".label"):
                continue                      # not/rozet etiketi: tasarim geregi
                                              # buyuk harf (KRITIK TESPIT, NOT,
                                              # YONTEM). Marka adi bir etikette
                                              # gectiginde varyant sayilmamali.
            t = plain(raw)
            if alinti_mi(t):
                continue                      # birebir alinti: model/kullanici
                                              # metni destenin kendi sesi degil
            quoted = [(q.start(), q.end()) for q in
                      re.finditer(r"[\"'\u201c\u201d\u2018\u2019][^\"'\u201c\u201d\u2018\u2019]{1,60}"
                                  r"[\"'\u201c\u201d\u2018\u2019]", t)]
            for m in re.finditer(re.escape(brand), t, re.IGNORECASE):
                if any(a <= m.start() < b for a, b in quoted):
                    continue              # literal filtre ifadesi / alinti
                tail = t[m.end():m.end() + 6]
                if re.match(r"\.(com|net|org|tr|ae|de|es|co)\b", tail):
                    continue                  # domain adi
                if re.search(r"[/\w.-]$", t[max(0, m.start() - 1):m.start()]):
                    continue                  # URL yolu icinde
                variants.add(m.group(0))
        if len(variants) > 1:
            rep.err("YAPI", "marka yazım tutarlılığı",
                    f"{sorted(variants)} (deste kendi sesinde; tablo hücreleri "
                    f"ve domain adları hariç)",
                    f"tüm destede tek yazım: '{brand}'")

    if slides and slides[-1].get("type") != "closing":
        rep.warn("YAPI", "kapanış", "son slayt 'closing' değil",
                 "deste 'Teşekkürler' slaytıyla kapanır")


# ----------------------------------------------------------------------------
# Katman 1: yerlesim (uretici olcum motoru)
# ----------------------------------------------------------------------------

def check_layout(spec, base, rep):
    here = os.path.dirname(os.path.abspath(__file__))
    assets = os.path.normpath(os.path.join(here, "..", "assets"))
    try:
        ctx, _n = build(spec, None, assets, base, check_only=True)
    except Exception as e:                       # pragma: no cover
        rep.err("YERLEŞİM", "üretim hatası", str(e), "deck.json şemasını kontrol et")
        return
    for w in ctx.warnings:
        if w.startswith("TASMA"):
            rep.err("YERLEŞİM", "taşma", w,
                    "blok yüksekliğini (h) düşür, font_pt küçült veya slaytı böl")
        elif "kucultuldu" in w:
            rep.warn("YERLEŞİM", "başlık auto-shrink", w,
                     "başlık tek satırda kaldı; kısaltmak daha temiz olabilir")
        else:
            rep.warn("YERLEŞİM", "ölçüm", w, "")


def beyaz_metin_denetimi(slide, no, prs, rep):
    """Beyaz metin, koyu dolgulu bir seklin uzerinde mi duruyor?

    Kart / panel / ayrac icindeki metinler beyazdir. Bir metin kutusu kabinin
    disina tastiginda PowerPoint'te fark edilmeyebilir ama beyaz zeminde
    gorunmez olur; Google Slides'a aktarildiginda satir tamamen kaybolur.
    Gercek olay: KPI kartinda 'h' kucultuldugunde delta satiri kartin altina
    tasti, PPTX'te beyaz uzerine beyaz kaldi (bkz. tuzaklar 3.6g).
    """
    from pptx.util import Emu
    PX = 9525.0

    def rect(sh):
        if None in (sh.left, sh.top, sh.width, sh.height):
            return None
        return (sh.left / PX, sh.top / PX,
                (sh.left + sh.width) / PX, (sh.top + sh.height) / PX)

    # Slayt zemini koyu renkliyse (ayrac, kapak, kapanis) tum yuzey kaptir.
    try:
        bg = slide.background.fill
        if bg.type is not None and int(bg.type) == 1 and str(bg.fore_color.rgb) != "FFFFFF":
            return
    except Exception:
        pass

    dolgu = []
    for sh in slide.shapes:
        r = rect(sh)
        if r is None:
            continue
        if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
            dolgu.append(r)               # tam sayfa kapak/ayrac gorseli
            continue
        try:
            if sh.fill.type is not None and int(sh.fill.type) == 1:      # solid
                c = sh.fill.fore_color
                if c.type is not None and str(c.rgb) != "FFFFFF":
                    dolgu.append(r)
        except Exception:
            continue

    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        beyaz = False
        for p in sh.text_frame.paragraphs:
            for run in p.runs:
                try:
                    if run.font.color.type is not None and str(run.font.color.rgb) == "FFFFFF":
                        beyaz = True
                except Exception:
                    pass
        if not beyaz:
            continue
        r = rect(sh)
        if r is None:
            continue
        # metin kutusu kendi yuksekligini asabildigi icin dikeyde tolerans
        icinde = any(d[0] - 2 <= r[0] and d[1] - 2 <= r[1]
                     and r[2] <= d[2] + 2 and r[3] <= d[3] + 6 for d in dolgu)
        if not icinde:
            rep.err(f"S{no:02d}", "beyaz metin dolgusuz alanda",
                    f"'{sh.text_frame.text.strip()[:44]}' kutusu hiçbir dolgulu "
                    f"şeklin içinde değil",
                    "kutuyu kabın içine al veya kabı büyüt; beyaz zeminde "
                    "görünmez olur ve Google Slides'ta tamamen kaybolur")


def check_pptx(path, rep):
    """Uretilmis dosyada bayat deger / yazim taramasi ve font denetimi."""
    try:
        from pptx import Presentation
    except Exception:
        return
    if not os.path.exists(path):
        rep.warn("DOSYA", "pptx", f"bulunamadı: {path}", "")
        return
    prs = Presentation(path)
    from pptx.util import Emu
    w, h = Emu(prs.slide_width).inches, Emu(prs.slide_height).inches
    if abs(w - 13.3333) > 0.02 or abs(h - 7.5) > 0.02:
        rep.err("DOSYA", "canvas", f"{w:.3f}x{h:.3f} inch",
                "13.333x7.5 inch olmalı (1280x720 px @96dpi)")
    bad_fonts = set()
    yanlis_display = []
    for i, s in enumerate(prs.slides, 1):
        for shp in s.shapes:
            if not shp.has_text_frame:
                continue
            for p in shp.text_frame.paragraphs:
                for r in p.runs:
                    n = r.font.name
                    if n and n not in ("Bricolage Grotesque", "Outfit"):
                        bad_fonts.add(n)
                    # Display kurali: 20pt ustu her metin Bricolage Grotesque.
                    # Kapak, ajanda, ayrac ve KPI degeri bu boydadir; govde
                    # metni hicbir yerde 20pt'ye cikmaz. Run duzeyinde yanlis
                    # aile yazildiginda kutu ayari eziliyordu (tuzaklar 3.6i).
                    try:
                        boy = r.font.size.pt if r.font.size else 0
                    except Exception:
                        boy = 0
                    if boy >= 20 and n == "Outfit" and r.text.strip():
                        yanlis_display.append(
                            f"S{i:02d} · {boy:.0f}pt · '{r.text.strip()[:32]}'")
    if yanlis_display:
        rep.err("DOSYA", "display fontu",
                f"{len(yanlis_display)} metin 20pt üstünde Outfit ile basılmış: "
                f"{yanlis_display[:4]}",
                "büyük puntolu metin Bricolage Grotesque olmalı")
        beyaz_metin_denetimi(s, i, prs, rep)
        if s.has_notes_slide:
            nt = s.notes_slide.notes_text_frame.text.lower()
            for ic in INTERNAL + TOOL_LEAK:
                if ic in nt:
                    rep.err(f"S{i:02d} NOT", "konuşmacı notu sızıntısı",
                            f"'{ic}'", "notları teslim öncesi temizle")
    if bad_fonts:
        rep.err("DOSYA", "font", f"{sorted(bad_fonts)}",
                "yalnızca Bricolage Grotesque ve Outfit kullanılır")


# ----------------------------------------------------------------------------

def main():
    ap = argparse.ArgumentParser(description="Inbound deste teslim denetimi")
    ap.add_argument("spec")
    ap.add_argument("--pptx", default=None)
    ap.add_argument("--json", action="store_true")
    ap.add_argument("--strict", action="store_true",
                    help="uyarilari da hata sayar")
    a = ap.parse_args()

    with open(a.spec, encoding="utf-8") as f:
        spec = json.load(f)
    base = os.path.dirname(os.path.abspath(a.spec))

    rep = Report()
    check_layout(spec, base, rep)
    check_language(spec, rep)
    check_numbers(spec, rep)
    check_structure(spec, rep)
    if a.pptx:
        check_pptx(a.pptx, rep)

    if a.json:
        print(json.dumps(rep.items, ensure_ascii=False, indent=1))
    else:
        if not rep.items:
            print("Denetim temiz: hata ve uyarı yok.")
        for lvl in ("HATA", "UYARI"):
            group = [i for i in rep.items if i["level"] == lvl]
            if not group:
                continue
            print(f"\n{'='*72}\n{lvl} ({len(group)})\n{'='*72}")
            for i in group:
                print(f"[{i['where']}] {i['rule']}: {i['detail']}")
                if i["fix"]:
                    print(f"    -> {i['fix']}")
        print(f"\nToplam: {len(rep.errors)} hata, {len(rep.warnings)} uyarı")

    if rep.errors or (a.strict and rep.warnings):
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
